from __future__ import annotations

import importlib.util
from pathlib import Path
import sys

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation


def _load_demo_module():
    repo_root = Path(__file__).resolve().parents[1]
    module_path = repo_root / "docs/figure/rr_peak_band_metric/plot_rr_peak_band_metric_demo.py"
    spec = importlib.util.spec_from_file_location("rr_peak_band_metric_demo", module_path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def test_editable_slide_contains_text_shapes_and_median_interval_explanation(tmp_path: Path):
    module = _load_demo_module()
    signal_png = tmp_path / "signal.png"
    pptx_path = tmp_path / "demo.pptx"

    fig, ax = plt.subplots(figsize=(4, 1.6))
    ax.plot(np.linspace(0, 1, 12), np.sin(np.linspace(0, 2 * np.pi, 12)))
    fig.savefig(signal_png)
    plt.close(fig)

    pred_details = module.PeakDetails(
        filtered=np.zeros(10),
        peaks=np.array([100, 650, 1200]),
        intervals_sec=np.array([5.5, 5.5]),
        median_interval_sec=5.5,
        rr_bpm=10.909,
        min_distance_samples=200,
        min_distance_sec=2.0,
        prominence=0.1,
        prominence_rule="demo",
    )
    target_details = module.PeakDetails(
        filtered=np.zeros(10),
        peaks=np.array([110, 670, 1230]),
        intervals_sec=np.array([5.6, 5.6]),
        median_interval_sec=5.6,
        rr_bpm=10.714,
        min_distance_samples=200,
        min_distance_sec=2.0,
        prominence=0.1,
        prominence_rule="demo",
    )

    module._write_editable_slide_pptx(
        signal_image_path=signal_png,
        pptx_path=pptx_path,
        dataset_row_id=5058,
        fs=100.0,
        low_hz=0.05,
        high_hz=0.70,
        order=4,
        pred_details=pred_details,
        target_details=target_details,
        abs_error=0.195,
    )

    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    texts = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))

    assert "带通峰值呼吸率绝对误差" in texts
    assert "相邻峰间隔" in texts
    assert "中位数" in texts
    assert "60 / 中位峰间隔" in texts
    assert sum(1 for shape in slide.shapes if shape.shape_type == 13) == 1
    assert len(slide.shapes) > 6


def test_signal_panel_contains_only_bandpassed_signal_with_peaks():
    module = _load_demo_module()
    time = np.linspace(0, 2 * np.pi, 100)
    pred_details = module.PeakDetails(
        filtered=np.sin(time),
        peaks=np.array([10, 50, 90]),
        intervals_sec=np.array([4.0, 4.0]),
        median_interval_sec=4.0,
        rr_bpm=15.0,
        min_distance_samples=200,
        min_distance_sec=2.0,
        prominence=0.1,
        prominence_rule="demo",
    )
    target_details = module.PeakDetails(
        filtered=np.cos(time),
        peaks=np.array([20, 60]),
        intervals_sec=np.array([4.0]),
        median_interval_sec=4.0,
        rr_bpm=15.0,
        min_distance_samples=200,
        min_distance_sec=2.0,
        prominence=0.1,
        prominence_rule="demo",
    )

    fig = module._build_signal_panel_figure(
        pred_details=pred_details,
        target_details=target_details,
        fs=10.0,
    )
    try:
        axis_titles = [ax.get_title() for ax in fig.axes]
        assert len(fig.axes) == 1
        assert "原始预测波形与真实波形" not in axis_titles
        assert "带通后信号与峰值位置" in axis_titles
    finally:
        plt.close(fig)
