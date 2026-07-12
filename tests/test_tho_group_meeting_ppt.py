from __future__ import annotations

import csv
import hashlib
import json
import os
from pathlib import Path
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches
import pytest


REPO_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = REPO_ROOT / "docs/stage_reports/20260708/组会汇报.pptx"
FINAL_DECK = REPO_ROOT / "docs/stage_reports/20260708/THO_research_v2_阶段进展_组会汇报.pptx"


def test_discussion_evidence_catalog_resolves_formal_configs_and_signal_sources():
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    catalog = build_evidence_catalog(REPO_ROOT)

    assert catalog.dataset_root.exists()
    assert catalog.dataset_index.exists()
    assert catalog.general_signal_npz.name == "f0_visual_sample_signals.npz"
    assert catalog.general_signal_npz.exists()
    assert catalog.general_sample_row_id == 8025
    assert (catalog.result_root / "g_series_test_eval_manifest.csv").exists()
    assert (catalog.harmonic_root / "model_metrics" / "model_stratified_metrics_summary.csv").exists()
    assert catalog.run_configs["g3_c_wide_8p0"].stft_win == 2000
    assert catalog.run_configs["g3_c_wide_8p0"].stft_hop == 250
    assert catalog.run_configs["g3_c_bandenergy"].stft_encoder_type == "bandenergy"
    assert all(config.path.name == "config.yaml" for config in catalog.run_configs.values())
    assert all("g_series_stft_input" in config.path.parts for config in catalog.run_configs.values())
    assert catalog.case_row_ids == (640, 873, 1353, 3584)


def test_signal_assets_stft_matches_model_centered_shapes():
    import numpy as np

    from scripts.tho_group_meeting_ppt.signal_figures import compute_stft_logmag

    signal = np.linspace(-1.0, 1.0, 18_000, dtype=np.float32)
    f0 = compute_stft_logmag(
        signal,
        sample_rate_hz=100.0,
        win_samples=3000,
        hop_samples=500,
        low_hz=0.05,
        high_hz=8.0,
        center=True,
    )
    wide = compute_stft_logmag(
        signal,
        sample_rate_hz=100.0,
        win_samples=2000,
        hop_samples=250,
        low_hz=0.05,
        high_hz=8.0,
        center=True,
    )

    assert f0.log_magnitude.shape == (239, 37)
    assert wide.log_magnitude.shape == (160, 73)
    assert f0.center is True and wide.center is True
    assert f0.frequencies_hz[0] == pytest.approx(1 / 15)
    assert wide.frequencies_hz[0] == pytest.approx(0.05)
    assert wide.frequency_resolution_hz == pytest.approx(0.05)


def test_signal_assets_rejects_signal_metadata_stft_mismatch(tmp_path: Path):
    from types import SimpleNamespace

    from scripts.tho_group_meeting_ppt.signal_figures import validate_signal_metadata_stft

    metadata_path = tmp_path / "f0_visual_sample_metadata.json"
    metadata = {
        "stft_win_samples": 2999,
        "stft_hop_samples": 500,
        "stft_low_hz": 0.05,
        "stft_high_hz": 8.0,
    }
    config = SimpleNamespace(stft_win=3000, stft_hop=500, stft_low_hz=0.05, stft_high_hz=8.0)

    with pytest.raises(
        ValueError,
        match=r"stft_win_samples.*2999.*3000.*f0_visual_sample_metadata\.json",
    ):
        validate_signal_metadata_stft(metadata, config, metadata_path)


def test_signal_assets_bandenergy_uses_repository_overlapping_inclusive_bands():
    import numpy as np

    from scripts.tho_group_meeting_ppt.signal_figures import (
        BANDENERGY_BANDS_HZ,
        bandenergy_from_logmag,
    )

    frequencies = np.arange(0.05, 8.0001, 0.05, dtype=np.float32)
    logmag = np.repeat(frequencies[:, None], 3, axis=1)
    energies = bandenergy_from_logmag(logmag, frequencies)

    assert BANDENERGY_BANDS_HZ == (
        (0.05, 0.3),
        (0.1, 0.7),
        (0.3, 1.2),
        (0.7, 3.0),
        (3.0, 8.0),
    )
    assert energies.shape == (5, 3)
    for idx, (low, high) in enumerate(BANDENERGY_BANDS_HZ):
        # rFFT 的 float32 频点会把 0.7 表示为 0.70000005；仓库实现按 bin
        # 索引纳入边界，因此断言也显式允许该表示误差。
        expected = frequencies[
            (frequencies >= low - 1e-6) & (frequencies <= high + 1e-6)
        ].mean()
        assert energies[idx] == pytest.approx(expected)


def test_signal_assets_footer_layout_stays_inside_and_clear_of_xlabels():
    import matplotlib.pyplot as plt

    from scripts.tho_group_meeting_ppt.signal_figures import (
        create_figure_with_footer,
        inspect_selected_text_bboxes,
    )

    fig, axes, footer = create_figure_with_footer(
        rows=2,
        columns=2,
        footer_text="Parameters come from provenance; this is a sufficiently long evidence note.",
        sharex=True,
    )
    for index, axis in enumerate(axes.flat):
        axis.plot([0, 1], [index, index + 1])
        axis.set_title(f"Panel {index + 1}")
        axis.set_xlabel("Time in window (s)")
        axis.set_ylabel("Amplitude")

    status = axes.flat[0].text(3.0, 0.5, "Critical status", transform=axes.flat[0].transAxes)
    outside_report = inspect_selected_text_bboxes(fig, footer, axes.flat)
    assert outside_report["all_key_text_inside"] is False
    status.set_position((0.5, 0.5))

    report = inspect_selected_text_bboxes(fig, footer, axes.flat)

    assert report["all_key_text_inside"] is True
    assert report["footer_overlaps_xlabels"] == 0
    plt.close(fig)


def test_signal_assets_bandenergy_labels_do_not_overlap_each_other_or_colorbar():
    import matplotlib.pyplot as plt
    import numpy as np
    import warnings

    from scripts.tho_group_meeting_ppt.signal_figures import (
        StftLogMagnitude,
        build_bandenergy_figure,
        inspect_selected_text_bboxes,
    )

    wide = StftLogMagnitude(
        log_magnitude=np.ones((160, 73), dtype=np.float32),
        frequencies_hz=np.arange(0.05, 8.0001, 0.05, dtype=np.float32),
        times_sec=np.arange(73, dtype=np.float32) * 2.5,
        frequency_resolution_hz=0.05,
        center=True,
        win_samples=2000,
        hop_samples=250,
    )
    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        fig = build_bandenergy_figure(wide)
        report = inspect_selected_text_bboxes(
            fig,
            getattr(fig, "_signal_footer_artist"),
            getattr(fig, "_signal_panel_axes"),
        )

    assert report["tagged_text_pair_overlaps"] == 0
    assert report["tagged_text_colorbar_overlaps"] == 0
    assert report["tagged_text_count"] == 5
    assert not [warning for warning in caught if "Glyph" in str(warning.message)]
    plt.close(fig)


def test_signal_assets_softz_explanation_limits_single_point_operator_claims():
    from scripts.tho_group_meeting_ppt.signal_figures import softz_explanation

    title, note = softz_explanation(
        knee=6.0,
        scale=2.0,
        ramp_s=5.0,
        max_abs_before=2.473841,
        changed_samples=0,
    )

    text = f"{title}\n{note}"
    assert "保留排序" not in text
    assert "full-weight" in text
    assert "raw + w(t)" in text
    assert "未触发压缩" in text
    assert "0 samples changed" in text


def test_signal_assets_preprocessing_gap_records_missing_fields_without_fake_curves(tmp_path: Path):
    from types import SimpleNamespace

    import numpy as np

    from scripts.tho_group_meeting_ppt.signal_figures import load_preprocessing_window

    dataset = tmp_path / "dataset"
    index = dataset / "training" / "dataset_index.csv"
    index.parent.mkdir(parents=True)
    source = dataset / "whole_night" / "source.npz"
    target = dataset / "whole_night" / "target.npz"
    source.parent.mkdir(parents=True)
    np.savez(source, unrelated=np.arange(30_000, dtype=np.float32))
    np.savez(target, unrelated=np.arange(30_000, dtype=np.float32))
    index.write_text(
        "dataset_row_id,window_start_s,window_end_s,source_npz,target_source_npz,"
        "bcg_rawish_observed_key,bcg_rawish_segment_robust_z_key,"
        "target_waveform_observed_key,target_waveform_segment_robust_z_key,"
        "target_waveform_segment_soft_z_key\n"
        "8025,90,270,../whole_night/source.npz,../whole_night/target.npz,"
        "bcg_raw,bcg_robust,tho_raw,tho_robust,tho_soft\n",
        encoding="utf-8",
    )
    catalog = SimpleNamespace(dataset_index=index)

    loaded = load_preprocessing_window(catalog, tmp_path / "out")
    gap = json.loads((tmp_path / "out" / "evidence_gap_preprocessing.json").read_text(encoding="utf-8"))

    assert loaded is None
    assert gap["row_id"] == 8025
    assert {
        "bcg_raw",
        "bcg_robust",
        "tho_raw",
        "tho_robust",
        "tho_soft",
    } <= set(gap["missing_fields"])
    assert any(field.startswith("soft-z parameters:") for field in gap["missing_fields"])
    assert gap["checked_npz_and_keys"]
    assert "suggested_generation_entrypoint" in gap


def test_signal_assets_build_keeps_missing_preprocessing_evidence_generatable(
    tmp_path: Path,
    monkeypatch,
):
    from types import SimpleNamespace

    import numpy as np

    from scripts.tho_group_meeting_ppt import signal_figures

    repo = tmp_path / "repo"
    signal_dir = repo / "evidence"
    signal_dir.mkdir(parents=True)
    signal_npz = signal_dir / "f0_visual_sample_signals.npz"
    time = np.arange(18_000, dtype=np.float32) / 100.0
    wave = np.sin(2 * np.pi * 0.25 * time).astype(np.float32)
    np.savez(
        signal_npz,
        time_sec_full=time,
        bcg_input_full=wave,
        target_respiration_full=wave,
        f0_prediction_full=wave,
    )
    source_checkpoint = signal_dir / "checkpoint.pt"
    source_run_config = signal_dir / "source_config.yaml"
    source_checkpoint.write_bytes(b"synthetic checkpoint evidence")
    source_run_config.write_text("model: {}\n", encoding="utf-8")
    signal_npz.with_name("f0_visual_sample_metadata.json").write_text(
        json.dumps(
            {
                "dataset_row_id": 8025,
                "crop_start_sec": 30.0,
                "crop_duration_sec": 60.0,
                "stft_win_samples": 3000,
                "stft_hop_samples": 500,
                "stft_low_hz": 0.05,
                "stft_high_hz": 8.0,
                "source_checkpoint": str(source_checkpoint.relative_to(repo)),
                "source_run_config": str(source_run_config.relative_to(repo)),
            }
        ),
        encoding="utf-8",
    )
    dataset = repo / "dataset"
    index = dataset / "training" / "dataset_index.csv"
    index.parent.mkdir(parents=True)
    index.write_text(
        "dataset_row_id,window_start_s,window_end_s,source_npz,target_source_npz,"
        "bcg_rawish_observed_key,bcg_rawish_segment_robust_z_key,"
        "target_waveform_observed_key,target_waveform_segment_robust_z_key,"
        "target_waveform_segment_soft_z_key\n"
        "8025,90,270,../whole_night/missing_source.npz,"
        "../whole_night/missing_target.npz,bcg_raw,bcg_robust,tho_raw,tho_robust,tho_soft\n",
        encoding="utf-8",
    )

    def config(label: str, *, win: int, hop: int, encoder: str):
        path = signal_dir / f"{label}.yaml"
        path.write_text("model: {}\n", encoding="utf-8")
        return SimpleNamespace(
            path=path,
            stft_win=win,
            stft_hop=hop,
            stft_low_hz=0.05,
            stft_high_hz=8.0,
            stft_encoder_type=encoder,
            stft_inject_position="pre_mixer",
        )

    catalog = SimpleNamespace(
        dataset_root=dataset,
        dataset_index=index,
        general_signal_npz=signal_npz,
        general_sample_row_id=8025,
        run_configs={
            "g0_f0_native_stft_pre_mixer": config("f0", win=3000, hop=500, encoder="conv2d"),
            "g3_c_wide_8p0": config("wide", win=2000, hop=250, encoder="conv2d"),
            "g3_c_bandenergy": config("bandenergy", win=2000, hop=250, encoder="bandenergy"),
        },
    )
    monkeypatch.setattr(signal_figures, "build_evidence_catalog", lambda _: catalog)
    gap_titles: list[str] = []
    original_gap_panel = signal_figures._gap_panel

    def recording_gap_panel(path: Path, title: str) -> None:
        gap_titles.append(title)
        original_gap_panel(path, title)

    monkeypatch.setattr(signal_figures, "_gap_panel", recording_gap_panel)
    output = tmp_path / "output"

    assets, metadata = signal_figures.build_signal_assets(repo, output)
    manifest = json.loads((output / "signal_assets_manifest.json").read_text(encoding="utf-8"))

    assert set(assets) == {
        "signal_overview",
        "preprocessing_comparison",
        "softz_mapping",
        "stft_resolution_comparison",
        "bandenergy_response",
    }
    assert all(path.is_file() for path in assets.values())
    assert metadata["preprocessing_evidence_gap"] is True
    assert set(gap_titles) == {"预处理前后对照", "soft-z 映射与真实分布"}
    assert (output / "evidence_gap_preprocessing.json").is_file()
    for key in ("source_npz", "target_npz"):
        record = manifest["evidence"][key]
        assert record["status"] == "missing"
        assert "sha256" not in record and "size_bytes" not in record


def test_signal_assets_builds_five_readable_real_evidence_figures(tmp_path: Path):
    import hashlib

    from PIL import Image, ImageStat

    from scripts.tho_group_meeting_ppt.signal_figures import build_signal_assets

    assets, metadata = build_signal_assets(REPO_ROOT, tmp_path)

    assert set(assets) == {
        "signal_overview",
        "preprocessing_comparison",
        "softz_mapping",
        "stft_resolution_comparison",
        "bandenergy_response",
    }
    assert metadata["dataset_row_id"] == 8025
    assert metadata["signal_length"] == 18_000
    assert metadata["f0_frames"] == 37
    assert metadata["wide_frames"] == 73
    assert metadata["wide_frequency_resolution_hz"] == pytest.approx(0.05)
    assert metadata["softz_changed_samples"] == 0
    assert metadata["softz_max_abs_before"] == pytest.approx(2.473841, abs=1e-6)
    manifest_path = tmp_path / "signal_assets_manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert manifest["dataset_row_id"] == 8025
    assert manifest["resolved_configs"]["f0"]["stft_win"] == 3000
    assert manifest["resolved_configs"]["wide"]["stft_hop"] == 250
    assert manifest["resolved_configs"]["bandenergy"]["stft_encoder_type"] == "bandenergy"
    assert manifest["evidence"]["dataset_index"]["path"]
    generator = manifest["evidence"]["signal_figure_code"]
    generator_path = REPO_ROOT / generator["path"]
    assert generator_path.name == "signal_figures.py"
    assert generator["sha256"] == hashlib.sha256(generator_path.read_bytes()).hexdigest()
    assert Path(manifest["evidence"]["source_npz"]["path"]).is_absolute()
    assert Path(manifest["evidence"]["target_npz"]["path"]).is_absolute()
    for path in assets.values():
        assert path.is_file() and path.stat().st_size > 20_000
        assert manifest["assets"][path.stem]["sha256"] == hashlib.sha256(path.read_bytes()).hexdigest()
        with Image.open(path) as image:
            assert image.width >= 1600 and image.height >= 900
            assert ImageStat.Stat(image.convert("L")).stddev[0] > 10


def test_read_run_config_evidence_reports_missing_field_with_path(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import read_run_config

    config = tmp_path / "config.yaml"
    config.write_text("model:\n  patch_len: 256\n", encoding="utf-8")

    with pytest.raises(ValueError, match=r"stft_win.*config\.yaml"):
        read_run_config(config)


def _evidence_config_text(*, dataset_root: str = "data/dataset", field_override: tuple[str, str] | None = None) -> str:
    model = {
        "patch_len": "256",
        "patch_stride": "128",
        "mixer_layers": "2",
        "base_channels": "16",
        "stft_win": "2000",
        "stft_hop": "250",
        "stft_low_hz": "0.05",
        "stft_high_hz": "8.0",
        "stft_encoder_type": "conv2d",
        "stft_inject_position": "pre_mixer",
    }
    if field_override is not None:
        model[field_override[0]] = field_override[1]
    model_yaml = "\n".join(f"  {key}: {value}" for key, value in model.items())
    return (
        "data:\n"
        f"  dataset_root: {dataset_root}\n"
        "  index_csv: training/dataset_index.csv\n"
        "  format: research_v2\n"
        "  input_set: research_v2_waveform\n"
        "  train_split: train\n"
        "  val_split: val\n"
        "  target_task: waveform\n"
        "  bcg_input_key: bcg_rawish_segment_soft_z_key\n"
        "  target_key: target_waveform_segment_soft_z_key\n"
        f"model:\n{model_yaml}\n"
    )


def _synthetic_evidence_repo(tmp_path: Path) -> tuple[Path, Path, Path, Path]:
    repo = tmp_path / "repo"
    dataset = repo / "data" / "dataset"
    (dataset / "training").mkdir(parents=True)
    (dataset / "training" / "dataset_index.csv").write_text("row_id\n1\n", encoding="utf-8")

    signal = repo / "assets" / "f0_visual_sample_signals.npz"
    signal.parent.mkdir(parents=True)
    signal.touch()
    signal.with_name("f0_visual_sample_metadata.json").write_text(
        json.dumps({"dataset_row_id": 8025, "files": {"signal_arrays_npz": signal.name}}),
        encoding="utf-8",
    )

    harmonic = repo / "harmonic"
    artifact_files = (
        harmonic / "test_v2" / "test_harmonic_labels.csv",
        harmonic / "model_metrics" / "model_stratified_metrics_summary.csv",
        harmonic / "corrections" / "model_harmonic_correction_summary.csv",
    )
    for path in artifact_files:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text("status\ncomplete\n", encoding="utf-8")

    labels = ("g0_time_only", "g0_f0_native_stft_pre_mixer", "g3_c_wide_8p0", "g3_c_bandenergy")
    seeds = (20260700, 20260837, 20260901)
    labels_path = harmonic / "test_v2" / "test_harmonic_labels.csv"
    shared_labels_hash = hashlib.sha256(labels_path.read_bytes()).hexdigest()
    thresholds_path = harmonic / "harmonic_thresholds.json"
    thresholds_path.write_text(json.dumps({"threshold_version": "synthetic-v1"}), encoding="utf-8")
    thresholds_hash = hashlib.sha256(thresholds_path.read_bytes()).hexdigest()
    (harmonic / "model_metrics" / "analysis_manifest.json").write_text(
        json.dumps({
            "operation": "summarize-metrics",
            "labels_path": str(labels_path.relative_to(repo)),
            "labels_sha256": shared_labels_hash,
            "eval_root": "results",
            "dataset_index": "data/dataset/training/dataset_index.csv",
            "model_labels": list(labels),
            "seeds": list(seeds),
        }),
        encoding="utf-8",
    )
    (harmonic / "test_v2" / "analysis_manifest.json").write_text(
        json.dumps({
            "operation": "apply",
            "dataset_root": "data/dataset",
            "index_csv": "training/dataset_index.csv",
            "split": "test",
            "thresholds_path": str(thresholds_path.relative_to(repo)),
            "thresholds_sha256": thresholds_hash,
            "threshold_version": "synthetic-v1",
            "outputs": {"labels": str(labels_path.relative_to(repo))},
        }),
        encoding="utf-8",
    )
    (harmonic / "corrections" / "analysis_manifest.json").write_text(
        json.dumps({
            "operation": "summarize-corrections",
            "labels_path": str(labels_path.relative_to(repo)),
            "labels_sha256": shared_labels_hash,
            "threshold_version": "synthetic-v1",
            "runs": [{"label": label, "seed": seed} for label in labels for seed in seeds],
        }),
        encoding="utf-8",
    )
    case_manifest = harmonic / "figures" / "model_case_manifest.csv"
    case_manifest.parent.mkdir(parents=True)
    case_manifest.write_text(
        "dataset_row_id,case_category\n"
        "640,all_corrected\n"
        "873,all_not_corrected\n"
        "1353,model_disagreement\n"
        "3584,threshold_boundary\n",
        encoding="utf-8",
    )

    manifest = repo / "results" / "g_series_test_eval_manifest.csv"
    manifest.parent.mkdir(parents=True)
    rows = []
    for label in labels:
        for seed in seeds:
            run_dir = repo / "runs" / label / str(seed)
            run_dir.mkdir(parents=True)
            (run_dir / "checkpoint_top1.pt").touch()
            (run_dir / "config.yaml").write_text(_evidence_config_text(), encoding="utf-8")
            output_stem = repo / "results" / f"{label}_{seed}"
            for suffix in ("_metrics.csv", "_summary.csv", "_manifest.csv"):
                Path(f"{output_stem}{suffix}").touch()
            rows.append({
                "label": label,
                "seed": str(seed),
                "checkpoint": str((run_dir / "checkpoint_top1.pt").relative_to(repo)),
                "metrics_output": str(Path(f"{output_stem}_metrics.csv").relative_to(repo)),
                "summary_output": str(Path(f"{output_stem}_summary.csv").relative_to(repo)),
                "manifest_output": str(Path(f"{output_stem}_manifest.csv").relative_to(repo)),
            })
    with manifest.open("w", newline="", encoding="utf-8") as fp:
        writer = csv.DictWriter(
            fp,
            fieldnames=("label", "seed", "checkpoint", "metrics_output", "summary_output", "manifest_output"),
        )
        writer.writeheader()
        writer.writerows(rows)
    return repo, manifest, harmonic, signal


def _enable_default_evidence_discovery(
    repo: Path,
    manifest: Path,
    harmonic: Path,
    signal: Path,
) -> tuple[Path, Path, Path]:
    discovered_manifest = repo / "runs" / "test_eval_g_series_synthetic_canonical" / manifest.name
    discovered_manifest.parent.mkdir(parents=True)
    manifest.rename(discovered_manifest)

    discovered_harmonic = repo / "runs" / "synthetic_harmonic_analysis"
    harmonic.rename(discovered_harmonic)
    relocated_labels = str(
        (discovered_harmonic / "test_v2" / "test_harmonic_labels.csv").relative_to(repo)
    )
    for relative in (
        Path("model_metrics/analysis_manifest.json"),
        Path("corrections/analysis_manifest.json"),
    ):
        provenance = discovered_harmonic / relative
        payload = json.loads(provenance.read_text(encoding="utf-8"))
        payload["labels_path"] = relocated_labels
        if relative == Path("model_metrics/analysis_manifest.json"):
            payload["eval_root"] = str(discovered_manifest.parent.relative_to(repo))
        provenance.write_text(json.dumps(payload), encoding="utf-8")
    test_manifest = discovered_harmonic / "test_v2" / "analysis_manifest.json"
    test_payload = json.loads(test_manifest.read_text(encoding="utf-8"))
    test_payload["outputs"]["labels"] = relocated_labels
    test_payload["thresholds_path"] = str(
        (discovered_harmonic / "harmonic_thresholds.json").relative_to(repo)
    )
    test_manifest.write_text(json.dumps(test_payload), encoding="utf-8")

    discovered_signal = repo / "docs" / "figure" / "synthetic" / signal.name
    discovered_signal.parent.mkdir(parents=True)
    signal.rename(discovered_signal)
    signal.with_name("f0_visual_sample_metadata.json").rename(
        discovered_signal.with_name("f0_visual_sample_metadata.json")
    )
    return discovered_manifest, discovered_harmonic, discovered_signal


def test_discussion_evidence_catalog_discovers_unique_structured_candidates(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    discovered = _enable_default_evidence_discovery(repo, manifest, harmonic, signal)

    catalog = build_evidence_catalog(repo)

    assert catalog.result_root == discovered[0].parent.resolve()
    assert catalog.harmonic_root == discovered[1].resolve()
    assert catalog.general_signal_npz == discovered[2].resolve()
    assert catalog.general_sample_row_id == 8025


def test_discussion_evidence_catalog_requires_explicit_manifest_when_candidates_are_ambiguous(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    discovered_manifest, _, _ = _enable_default_evidence_discovery(repo, manifest, harmonic, signal)
    second = repo / "runs" / "test_eval_g_series_second_canonical" / discovered_manifest.name
    second.parent.mkdir(parents=True)
    second.write_text(discovered_manifest.read_text(encoding="utf-8"), encoding="utf-8")

    with pytest.raises(
        ValueError,
        match=r"manifest.*候选.*(?:synthetic_canonical.*second_canonical|second_canonical.*synthetic_canonical).*显式",
    ):
        build_evidence_catalog(repo)


def test_discussion_evidence_catalog_accepts_explicit_paths_relative_to_repo_root(tmp_path: Path, monkeypatch):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    elsewhere = tmp_path / "elsewhere"
    elsewhere.mkdir()
    monkeypatch.chdir(elsewhere)

    catalog = build_evidence_catalog(
        repo,
        manifest_path=manifest.relative_to(repo),
        harmonic_root=harmonic.relative_to(repo),
        general_signal_npz=signal.relative_to(repo),
    )

    assert catalog.dataset_root == (repo / "data" / "dataset").resolve()
    assert catalog.dataset_index == (repo / "data" / "dataset" / "training" / "dataset_index.csv").resolve()
    assert catalog.result_root == manifest.parent.resolve()
    assert catalog.harmonic_root == harmonic.resolve()
    assert catalog.general_signal_npz == signal.resolve()


def test_evidence_manifest_rejects_duplicate_label_seed(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    with manifest.open(newline="", encoding="utf-8") as fp:
        rows = list(csv.DictReader(fp))
    rows.append(rows[0].copy())
    with manifest.open("w", newline="", encoding="utf-8") as fp:
        writer = csv.DictWriter(fp, fieldnames=rows[0].keys())
        writer.writeheader()
        writer.writerows(rows)

    with pytest.raises(ValueError, match=r"重复.*g0_time_only.*20260700.*g_series_test_eval_manifest\.csv"):
        build_evidence_catalog(repo, manifest_path=manifest, harmonic_root=harmonic, general_signal_npz=signal)


def test_evidence_manifest_checks_nonfirst_seed_data_provenance(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    mismatched = repo / "runs" / "g3_c_wide_8p0" / "20260837" / "config.yaml"
    other_dataset = repo / "data" / "other"
    (other_dataset / "training").mkdir(parents=True)
    (other_dataset / "training" / "dataset_index.csv").touch()
    mismatched.write_text(_evidence_config_text(dataset_root="data/other"), encoding="utf-8")

    with pytest.raises(
        ValueError,
        match=r"g3_c_wide_8p0.*20260837.*data\.dataset_root.*config\.yaml",
    ):
        build_evidence_catalog(repo, manifest_path=manifest, harmonic_root=harmonic, general_signal_npz=signal)


def test_evidence_harmonic_provenance_rejects_mismatched_eval_root(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    provenance = harmonic / "model_metrics" / "analysis_manifest.json"
    payload = json.loads(provenance.read_text(encoding="utf-8"))
    payload["eval_root"] = "results/wrong_eval"
    provenance.write_text(json.dumps(payload), encoding="utf-8")

    with pytest.raises(
        ValueError,
        match=r"model_metrics/analysis_manifest\.json.*eval_root.*wrong_eval.*results",
    ):
        build_evidence_catalog(repo, manifest_path=manifest, harmonic_root=harmonic, general_signal_npz=signal)


def test_evidence_harmonic_provenance_rejects_wrong_case_category(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    case_manifest = harmonic / "figures" / "model_case_manifest.csv"
    case_manifest.write_text(
        case_manifest.read_text(encoding="utf-8").replace("640,all_corrected", "640,wrong_category"),
        encoding="utf-8",
    )

    with pytest.raises(
        ValueError,
        match=r"model_case_manifest\.csv.*dataset_row_id=640.*case_category.*wrong_category.*all_corrected",
    ):
        build_evidence_catalog(repo, manifest_path=manifest, harmonic_root=harmonic, general_signal_npz=signal)


def test_evidence_harmonic_provenance_rejects_tampered_labels_hash(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    labels = harmonic / "test_v2" / "test_harmonic_labels.csv"
    labels.write_text(labels.read_text(encoding="utf-8") + "tampered\n", encoding="utf-8")

    with pytest.raises(
        ValueError,
        match=r"test_harmonic_labels\.csv.*labels_sha256.*实际.*期望",
    ):
        build_evidence_catalog(repo, manifest_path=manifest, harmonic_root=harmonic, general_signal_npz=signal)


def test_evidence_harmonic_provenance_rejects_tampered_thresholds_hash(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    repo, manifest, harmonic, signal = _synthetic_evidence_repo(tmp_path)
    thresholds = harmonic / "harmonic_thresholds.json"
    thresholds.write_text('{"tampered": true}', encoding="utf-8")

    with pytest.raises(
        ValueError,
        match=r"harmonic_thresholds\.json.*thresholds_sha256.*实际.*期望",
    ):
        build_evidence_catalog(repo, manifest_path=manifest, harmonic_root=harmonic, general_signal_npz=signal)


def _write_case_prediction(
    path: Path,
    *,
    label: str,
    row_ids: list[int],
    seed: int = 20260837,
) -> None:
    import numpy as np

    values = np.stack(
        [np.full(18_000, row_id, dtype=np.float32) for row_id in row_ids]
    )
    np.savez(
        path,
        dataset_row_id=np.asarray(row_ids, dtype=np.int64),
        r_tho_hat=values,
        tho_ref=values / 10.0,
    )
    path.with_name(f"{path.stem}_manifest.json").write_text(
        json.dumps(
            {
                "label": label,
                "seed": seed,
                "n_windows": len(row_ids),
                "array_schema": {
                    "dataset_row_id": [len(row_ids)],
                    "r_tho_hat": [len(row_ids), 18_000],
                    "tho_ref": [len(row_ids), 18_000],
                },
            }
        ),
        encoding="utf-8",
    )


def test_case_assets_prediction_join_is_by_dataset_row_id_not_npz_order(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.case_figures import load_case_predictions

    paths = []
    labels = ("time", "f0", "wide", "bandenergy")
    orders = ([30, 10, 20], [20, 30, 10], [10, 20, 30], [30, 20, 10])
    for label, order in zip(labels, orders, strict=True):
        path = tmp_path / f"{label}_20260837_harmonic_predictions.npz"
        _write_case_prediction(path, label=label, row_ids=list(order))
        paths.append(path)

    joined = load_case_predictions(paths, (10, 20, 30))

    assert set(joined) == set(labels)
    for label in labels:
        assert joined[label][20]["prediction"].shape == (18_000,)
        assert joined[label][20]["prediction"][0] == pytest.approx(20.0)
        assert joined[label][20]["prediction"].dtype.name == "float32"
        assert joined[label][20]["prediction"].base is None
        assert joined[label][20]["reference"].base is None


def test_case_assets_prediction_join_rejects_same_row_with_different_reference(tmp_path: Path):
    import numpy as np

    from scripts.tho_group_meeting_ppt.case_figures import load_case_predictions

    paths = []
    for label in ("time", "wide"):
        path = tmp_path / f"{label}_20260837_harmonic_predictions.npz"
        _write_case_prediction(path, label=label, row_ids=[10, 20, 30])
        paths.append(path)
    with np.load(paths[1]) as archive:
        prediction = np.asarray(archive["r_tho_hat"])
        reference = np.asarray(archive["tho_ref"]).copy()
        ids = np.asarray(archive["dataset_row_id"])
    reference[1, 0] += 1.0
    np.savez(paths[1], dataset_row_id=ids, r_tho_hat=prediction, tho_ref=reference)

    with pytest.raises(ValueError, match=r"row 20.*reference.*wide.*time"):
        load_case_predictions(paths, (10, 20, 30))


@pytest.mark.parametrize(("source", "field"), (("case_manifest", "samp_id"), ("correction:wide", "stratum")))
def test_case_assets_identity_rejects_subject_or_stratum_drift(source: str, field: str):
    from scripts.tho_group_meeting_ppt.case_figures import validate_case_identity

    records = {
        "dataset_index": {"dataset_row_id": "640", "samp_id": "220", "split": "test"},
        "harmonic_labels": {
            "dataset_row_id": "640",
            "samp_id": "220",
            "split": "test",
            "stratum": "harmonic_prominent",
        },
        "case_manifest": {"dataset_row_id": "640", "samp_id": "220", "seed": "20260837"},
        "correction:wide": {
            "dataset_row_id": "640",
            "samp_id": "220",
            "seed": "20260837",
            "input_stratum": "harmonic_prominent",
        },
    }
    records[source]["samp_id" if field == "samp_id" else "input_stratum"] = (
        "999" if field == "samp_id" else "peak_doubling"
    )

    with pytest.raises(ValueError, match=rf"row 640.*{source}.*{field}"):
        validate_case_identity(640, records, expected_seed=20260837)


def _write_prediction_provenance_fixture(tmp_path: Path) -> tuple[Path, Path, dict[str, str]]:
    import torch
    from omegaconf import OmegaConf

    run = tmp_path / "runs/model/run"
    run.mkdir(parents=True)
    checkpoint = run / "checkpoint_top1.pt"
    checkpoint_config = {"training": {"seed": 20260837}, "model": {"name": "fixture"}}
    torch.save({"config": checkpoint_config, "epoch": 1}, checkpoint)
    config = run / "config.yaml"
    OmegaConf.save(OmegaConf.create(checkpoint_config), config)
    labels = tmp_path / "labels.csv"
    labels.write_text("dataset_row_id,samp_id\n10,220\n", encoding="utf-8")
    prediction = tmp_path / "time_20260837_harmonic_predictions.npz"
    _write_case_prediction(prediction, label="time", row_ids=[10])
    eval_manifest = tmp_path / "eval_manifest.csv"
    eval_manifest.write_text(
        "checkpoint,config,split,sample_seed\n"
        f"{checkpoint.relative_to(tmp_path)},{config.relative_to(tmp_path)},test,20260837\n",
        encoding="utf-8",
    )
    payload = json.loads(prediction.with_name(f"{prediction.stem}_manifest.json").read_text())
    payload.update(
        {
            "split": "test",
            "checkpoint": str(checkpoint.relative_to(tmp_path)),
            "checkpoint_sha256": hashlib.sha256(checkpoint.read_bytes()).hexdigest(),
            "config": str(config.relative_to(tmp_path)),
            "labels_path": str(labels.relative_to(tmp_path)),
            "labels_sha256": hashlib.sha256(labels.read_bytes()).hexdigest(),
            "output_path": str(prediction.relative_to(tmp_path)),
        }
    )
    prediction.with_name(f"{prediction.stem}_manifest.json").write_text(json.dumps(payload), encoding="utf-8")
    canonical = {
        "label": "time",
        "seed": "20260837",
        "checkpoint": str(checkpoint.relative_to(tmp_path)),
        "manifest_output": str(eval_manifest.relative_to(tmp_path)),
    }
    return prediction, labels, canonical


@pytest.mark.parametrize(
    "tamper",
    ("missing_checkpoint_hash", "bad_labels_hash", "bad_config_hash", "config_outside_run", "wrong_config_seed"),
)
def test_case_prediction_provenance_rejects_missing_hash_or_tampering(tmp_path: Path, tamper: str):
    from omegaconf import OmegaConf

    from scripts.tho_group_meeting_ppt.case_figures import _validate_prediction_provenance

    prediction, labels, canonical = _write_prediction_provenance_fixture(tmp_path)
    manifest = prediction.with_name(f"{prediction.stem}_manifest.json")
    payload = json.loads(manifest.read_text())
    if tamper == "missing_checkpoint_hash":
        payload.pop("checkpoint_sha256")
    elif tamper == "bad_labels_hash":
        payload["labels_sha256"] = "0" * 64
    elif tamper == "bad_config_hash":
        payload["config_sha256"] = "0" * 64
    elif tamper == "config_outside_run":
        outside = tmp_path / "other.yaml"
        outside.write_text("training:\n  seed: 20260837\n", encoding="utf-8")
        payload["config"] = str(outside.relative_to(tmp_path))
    else:
        config = tmp_path / payload["config"]
        OmegaConf.save(OmegaConf.create({"training": {"seed": 1}}), config)
    manifest.write_text(json.dumps(payload), encoding="utf-8")

    with pytest.raises(ValueError, match=r"time.*(hash|config|seed).*(manifest|yaml|labels|checkpoint)"):
        _validate_prediction_provenance(tmp_path, prediction, canonical, labels)


def test_case_prediction_provenance_rejects_disk_config_drift_against_checkpoint(tmp_path: Path):
    from omegaconf import OmegaConf

    from scripts.tho_group_meeting_ppt.case_figures import _validate_prediction_provenance

    prediction, labels, canonical = _write_prediction_provenance_fixture(tmp_path)
    manifest = json.loads(
        prediction.with_name(f"{prediction.stem}_manifest.json").read_text(encoding="utf-8")
    )
    config = tmp_path / manifest["config"]
    drift = OmegaConf.load(config)
    OmegaConf.update(drift, "model.name", "drifted", merge=False)
    OmegaConf.save(drift, config)

    with pytest.raises(ValueError, match=r"time.*checkpoint.*embedded config.*config\.yaml"):
        _validate_prediction_provenance(tmp_path, prediction, canonical, labels)


def test_case_prediction_provenance_records_gap_when_no_config_hash_or_checkpoint_config(tmp_path: Path):
    import torch

    from scripts.tho_group_meeting_ppt.case_figures import _validate_prediction_provenance

    prediction, labels, canonical = _write_prediction_provenance_fixture(tmp_path)
    manifest_path = prediction.with_name(f"{prediction.stem}_manifest.json")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    checkpoint = tmp_path / manifest["checkpoint"]
    torch.save({"epoch": 1}, checkpoint)
    manifest["checkpoint_sha256"] = hashlib.sha256(checkpoint.read_bytes()).hexdigest()
    manifest.pop("config_sha256", None)
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

    _, evidence = _validate_prediction_provenance(tmp_path, prediction, canonical, labels)

    assert evidence["config_evidence_status"] == "path_and_seed_crosscheck_only"
    assert "缺少生成时 config hash" in evidence["config_evidence_gap"]
    assert "不能证明内容未漂移" in evidence["config_evidence_gap"]


@pytest.mark.parametrize(
    ("row_ids", "message"),
    (([10, 10, 20], r"time.*row 10.*predictions\.npz"), ([10, 20], r"time.*row 30.*predictions\.npz")),
)
def test_case_assets_prediction_join_rejects_duplicate_or_missing_row(
    tmp_path: Path,
    row_ids: list[int],
    message: str,
):
    from scripts.tho_group_meeting_ppt.case_figures import load_case_predictions

    path = tmp_path / "time_20260837_harmonic_predictions.npz"
    _write_case_prediction(path, label="time", row_ids=row_ids)

    with pytest.raises(ValueError, match=message):
        load_case_predictions([path], (10, 20, 30))


def test_case_assets_builds_real_delta_stability_strata_and_complete_cases(tmp_path: Path):
    from PIL import Image

    from scripts.tho_group_meeting_ppt.case_figures import build_case_assets

    assets, metadata = build_case_assets(REPO_ROOT, tmp_path)

    assert set(assets) == {
        "overall_delta",
        "seed_subject_stability",
        "strata_tradeoffs",
        "case_row_640",
        "case_row_873",
        "case_row_1353",
        "case_row_3584",
    }
    assert metadata["case_row_ids"] == [640, 873, 1353, 3584]
    assert metadata["models_per_case"] == 4
    assert metadata["signal_length"] == 18_000
    assert metadata["stability"]["uncertainty_unit"] == "seed"
    assert metadata["stability"]["window_overlap_used_as_independent_ci"] is False
    label_summary = {
        row["label"]: row
        for row in csv.DictReader(
            (REPO_ROOT / "runs/test_eval_g_series_20260709_local_rr_canonical/g_series_local_rr_canonical_label_summary.csv").open()
        )
    }
    recomputed_wide_delta = float(label_summary["g3_c_wide_8p0"]["rr_peak_band_robust_abs_error_mean_mean"]) - float(
        label_summary["g0_time_only"]["rr_peak_band_robust_abs_error_mean_mean"]
    )
    paired_rows = list(
        csv.DictReader(
            (REPO_ROOT / "runs/bcg_second_harmonic_20260710/model_metrics/paired_delta_vs_time_summary.csv").open()
        )
    )
    recomputed_count_delta = float(
        next(
            row
            for row in paired_rows
            if row["comparison"] == "bandenergy_vs_time" and row["stratum"] == "all_windows"
        )["delta_breath_count_zero_cross_bpm_error_mean"]
    )
    assert metadata["deltas"]["wide_vs_time"]["robust_rr_bpm"] == pytest.approx(recomputed_wide_delta)
    assert metadata["deltas"]["bandenergy_vs_time"]["count_bpm"] == pytest.approx(recomputed_count_delta)
    real_cases = {
        row["dataset_row_id"]: row["case_category"]
        for row in csv.DictReader(
            (REPO_ROOT / "runs/bcg_second_harmonic_20260710/figures/model_case_manifest.csv").open()
        )
        if int(row["dataset_row_id"]) in {640, 873, 1353, 3584}
    }
    assert metadata["case_categories"] == real_cases
    assert set(metadata["case_titles"]) == set(real_cases)
    manifest = json.loads((tmp_path / "case_assets_manifest.json").read_text(encoding="utf-8"))
    assert manifest["generator"] == "scripts.tho_group_meeting_ppt.case_figures.build_case_assets"
    assert len(manifest["sources"]) >= 20
    generator = REPO_ROOT / "scripts/tho_group_meeting_ppt/case_figures.py"
    generator_record = next(
        record for record in manifest["sources"] if record["path"].endswith("scripts/tho_group_meeting_ppt/case_figures.py")
    )
    assert generator_record["sha256"] == hashlib.sha256(generator.read_bytes()).hexdigest()
    assert generator_record["size_bytes"] == generator.stat().st_size
    assert metadata["correction_thresholds"]["correction_ratio_drop_min"] == pytest.approx(0.2)
    assert metadata["case_threshold_evidence"]["3584"]["boundary_distance"] == pytest.approx(
        abs(0.06433105468749982 / 0.1 - 1.0)
        + abs(0.5015774167611345 / 0.5392740654828524 - 1.0)
        + abs(0.17347360905723672 / 0.16474917204045256 - 1.0)
    )
    for row_id in ("640", "873", "1353", "3584"):
        model_rows = metadata["case_threshold_evidence"][row_id]["models"]
        assert set(model_rows) == {
            "g0_time_only",
            "g0_f0_native_stft_pre_mixer",
            "g3_c_wide_8p0",
            "g3_c_bandenergy",
        }
        assert all("harmonic_ratio_relative_drop" in row for row in model_rows.values())
        assert all("correction_condition_met" in row for row in model_rows.values())
    source_names = [record["path"] for record in manifest["sources"]]
    assert any(name.endswith("checkpoint_top1.pt") for name in source_names)
    assert any(name.endswith("config.yaml") for name in source_names)
    assert any(name.endswith("harmonic_predictions_manifest.json") for name in source_names)
    assert any(name.endswith("harmonic_predictions.npz") for name in source_names)
    assert metadata["prediction_config_provenance"]
    assert all(
        item["config_evidence_status"] == "checkpoint_embedded_config_match"
        for item in metadata["prediction_config_provenance"].values()
    )
    assert metadata["config_evidence_gaps"] == {}
    assert metadata["stability"]["subject_aggregation_source"]["scope"] == "repo"
    assert not Path(metadata["stability"]["subject_aggregation_source"]["path"]).is_absolute()
    for item in metadata["prediction_config_provenance"].values():
        assert item["prediction_manifest_config_sha256_present"] is False
        assert item["checkpoint_embedded_config_present"] is True
        assert item["config"]["scope"] == "repo"
        assert item["checkpoint"]["scope"] == "repo"
        assert not Path(item["config"]["path"]).is_absolute()
        assert not Path(item["checkpoint"]["path"]).is_absolute()
    for record in [*manifest["sources"], *manifest["assets"].values()]:
        if record["scope"] == "repo":
            assert not Path(record["path"]).is_absolute()
            assert ".." not in Path(record["path"]).parts
        else:
            assert record["scope"] == "external"
            assert Path(record["path"]).is_absolute()
    for key, path in assets.items():
        assert path.name == f"{key}.png"
        assert manifest["assets"][key]["sha256"] == hashlib.sha256(path.read_bytes()).hexdigest()
        with Image.open(path) as image:
            assert image.width >= 1800 and image.height >= 1000
            assert image.convert("L").getbbox() is not None
    assert all(length == 18_000 for length in metadata["case_prediction_lengths"].values())


def test_case_manifest_repo_records_resolve_after_checkout_root_moves(tmp_path: Path):
    import shutil

    from scripts.tho_group_meeting_ppt.case_figures import resolve_manifest_record

    manifest_path = (
        REPO_ROOT
        / "docs/stage_reports/20260708/generated_assets/discussion/case_assets_manifest.json"
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    records = [
        next(record for record in manifest["sources"] if record["path"].endswith("case_figures.py")),
        manifest["assets"]["case_row_3584"],
    ]
    moved_root = tmp_path / "moved-checkout"
    for record in records:
        source = resolve_manifest_record(record, REPO_ROOT)
        moved = moved_root / record["path"]
        moved.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, moved)
        resolved = resolve_manifest_record(record, moved_root)
        assert resolved == moved.resolve()
        assert hashlib.sha256(resolved.read_bytes()).hexdigest() == record["sha256"]


@pytest.mark.parametrize(
    ("field", "yaml_value"),
    (
        ("patch_len", "true"),
        ("patch_len", "2.9"),
        ("stft_low_hz", ".nan"),
        ("stft_encoder_type", "[conv2d]"),
    ),
)
def test_read_run_config_evidence_rejects_invalid_model_field_types(
    tmp_path: Path,
    field: str,
    yaml_value: str,
):
    from scripts.tho_group_meeting_ppt.evidence import read_run_config

    config = tmp_path / "config.yaml"
    config.write_text(_evidence_config_text(field_override=(field, yaml_value)), encoding="utf-8")

    with pytest.raises(ValueError, match=rf"model\.{field}.*config\.yaml"):
        read_run_config(config)


def test_body_text_is_black_and_table_has_only_three_horizontal_rules():
    from scripts.tho_group_meeting_ppt.theme import (
        BODY_BLACK,
        WHITE,
        add_body_text,
        add_three_line_table,
        is_three_line_table,
        new_content_slide,
    )

    prs = Presentation()
    slide = new_content_slide(prs, "测试标题", page_number=1)
    box = add_body_text(
        slide,
        ["正文内容"],
        x=Inches(1.0),
        y=Inches(1.5),
        width=Inches(5.0),
        height=Inches(1.0),
    )
    table = add_three_line_table(
        slide,
        ["方案", "指标"],
        [["A", "0.1"]],
        x=Inches(1.0),
        y=Inches(2.8),
        width=Inches(6.0),
        height=Inches(1.2),
    )

    assert box.text_frame.paragraphs[0].runs[0].font.color.rgb == BODY_BLACK
    assert is_three_line_table(table)
    assert all(cell.fill.fore_color.rgb == WHITE for row in table.rows for cell in row.cells)
    assert not table.first_row
    assert not table.first_col
    assert not table.horz_banding
    assert not table.vert_banding
    header_tags = [child.tag.rsplit("}", 1)[-1] for child in table.cell(0, 0)._tc.get_or_add_tcPr()]
    assert header_tags[:3] == ["lnT", "lnB", "solidFill"]


def test_main_deck_has_exactly_25_ordered_slides_and_numeric_sources():
    from scripts.tho_group_meeting_ppt.content import MAIN_SLIDES

    assert len(MAIN_SLIDES) == 25
    assert MAIN_SLIDES[0].key == "title"
    assert MAIN_SLIDES[10].key == "overall_test_results"
    assert MAIN_SLIDES[24].key == "next_stage_takeaways"
    assert all(slide.sources for slide in MAIN_SLIDES if slide.contains_numeric_evidence)
    assert MAIN_SLIDES[0].placeholder == "【待补：汇报人、汇报日期】"


EXPECTED_DISCUSSION_SECTIONS = {
    "任务与信号直觉",
    "数据来源与样本形成",
    "输入与目标预处理",
    "模型输入与计算图",
    "训练与损失",
    "指标计算与失效场景",
    "对照实验设计",
    "整体结果与稳定性",
    "分层分析与完整案例",
    "BCG 二次谐波问题",
    "研究议题与下一步",
}


def test_discussion_units_cover_complete_causal_chain():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    assert {unit.section for unit in DISCUSSION_UNITS} == EXPECTED_DISCUSSION_SECTIONS
    assert len(DISCUSSION_UNITS) >= 42
    assert len({unit.key for unit in DISCUSSION_UNITS}) == len(DISCUSSION_UNITS)


def test_discussion_unit_schema_is_frozen_and_uses_tuples():
    from dataclasses import FrozenInstanceError, fields
    from typing import get_type_hints

    from scripts.tho_group_meeting_ppt.discussion_content import DiscussionUnit

    required_fields = {
        "key", "section", "title", "kind", "question", "method_steps", "parameters",
        "rationale", "evidence", "limits", "discussion_prompt", "visual_keys", "sources",
    }
    assert required_fields <= {field.name for field in fields(DiscussionUnit)}
    hints = get_type_hints(DiscussionUnit)
    for name in ("method_steps", "parameters", "rationale", "evidence", "limits", "visual_keys", "sources"):
        assert hints[name] == tuple[str, ...]
    assert hints["discussion_prompt"] == str | None
    unit = DiscussionUnit(key="k", section="s", title="t", kind="technical", question="q")
    with pytest.raises(FrozenInstanceError):
        unit.title = "changed"
    for name in ("method_steps", "parameters", "rationale", "evidence", "limits", "visual_keys", "sources"):
        assert isinstance(getattr(unit, name), tuple)


def test_discussion_technical_units_are_actionable_and_traceable():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    allowed_kinds = {"title", "section", "technical"}
    assert {unit.kind for unit in DISCUSSION_UNITS} <= allowed_kinds
    technical = [unit for unit in DISCUSSION_UNITS if unit.kind == "technical"]
    assert all(unit.method_steps or unit.visual_keys for unit in technical)
    assert all(unit.sources for unit in technical)


def test_discussion_units_have_no_blank_scalar_or_tuple_content():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    scalar_fields = ("key", "section", "title", "kind", "question")
    tuple_fields = ("method_steps", "parameters", "rationale", "evidence", "limits", "visual_keys", "sources")
    for unit in DISCUSSION_UNITS:
        assert all(isinstance(getattr(unit, name), str) and getattr(unit, name).strip() for name in scalar_fields), unit.key
        for name in tuple_fields:
            assert all(isinstance(value, str) and value.strip() for value in getattr(unit, name)), (unit.key, name)
        if unit.discussion_prompt is not None:
            assert isinstance(unit.discussion_prompt, str) and unit.discussion_prompt.strip(), unit.key


def test_high_risk_discussion_units_encode_required_domain_semantics():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    by_key = {unit.key: unit for unit in DISCUSSION_UNITS}
    patch = "\n".join((*by_key["patchmixer_token_shape"].method_steps, *by_key["patchmixer_token_shape"].parameters))
    assert "patch_len" in patch and "stride" in patch and "token" in patch

    loss = "\n".join((*by_key["composite_loss_goal"].method_steps, *by_key["composite_loss_goal"].rationale))
    assert "包络" in loss and "频谱" in loss and "方向" in loss

    robust_rr = "\n".join((*by_key["metric_robust_rr"].method_steps, *by_key["metric_robust_rr"].limits))
    assert "Welch" in robust_rr and "峰间距" in robust_rr and "双峰" in robust_rr

    paired = "\n".join((*by_key["paired_delta"].method_steps, *by_key["paired_delta"].limits))
    seed = "\n".join((*by_key["seed_stability"].method_steps, *by_key["seed_stability"].limits))
    assert "配对" in paired and "三个 seed" in paired
    assert "seed" in seed and "样本量小" in seed

    correction = "\n".join((*by_key["harmonic_correction"].method_steps, *by_key["harmonic_correction"].limits))
    assert "10%" in correction and "20%" in correction and "高纠正率不等于" in correction


def test_named_high_value_discussion_units_have_complete_decision_context():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    by_key = {unit.key: unit for unit in DISCUSSION_UNITS}

    def assert_complete(key: str) -> None:
        unit = by_key[key]
        for field_name in ("method_steps", "parameters", "rationale", "limits"):
            values = getattr(unit, field_name)
            assert values, (key, field_name)
            assert all(isinstance(value, str) and value.strip() for value in values), (key, field_name)
        assert unit.discussion_prompt is not None, key
        assert isinstance(unit.discussion_prompt, str) and unit.discussion_prompt.strip(), key

    for key in (
        "target_soft_z_and_mask",
        "patchmixer_token_shape",
        "composite_loss_goal",
        "metric_robust_rr",
        "paired_delta",
        "harmonic_correction",
    ):
        assert_complete(key)


def test_target_supervision_unit_distinguishes_window_filter_loss_and_rr_mask():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    unit = next(unit for unit in DISCUSSION_UNITS if unit.key == "target_soft_z_and_mask")
    text = "\n".join((*unit.method_steps, *unit.parameters, *unit.rationale, *unit.evidence, *unit.limits))
    assert "整窗筛选" in text
    assert "完整 target 波形" in text
    assert "rr_peak_valid_mask" in text
    assert "RR" in text and "评价" in text
    assert "仅在允许位置形成监督" not in text


def test_overall_count_result_states_cycle_and_bpm_units():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    unit = next(unit for unit in DISCUSSION_UNITS if unit.key == "overall_metric_values")
    parameters = "\n".join(unit.parameters)
    assert "0.712186 bpm" in parameters
    assert "0.870774" in parameters and "无量纲" in parameters
    assert "2.054401" in parameters and "周期数绝对误差" in parameters
    assert "0.684800" in parameters and "bpm" in parameters
    assert "0.773299 bpm" in parameters


def test_checkpoint_selection_describes_actual_g_series_topk_and_reevaluation():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    unit = next(unit for unit in DISCUSSION_UNITS if unit.key == "checkpoint_selection")
    text = "\n".join((*unit.method_steps, *unit.parameters, *unit.evidence, *unit.limits))
    assert "方向 gate" in text
    assert "val-loss top-k" in text
    assert "旁路复评" in text
    assert "runs/test_eval_g_series_20260709_local_rr_canonical/g_series_test_eval_manifest.csv" in text
    assert "checkpoint_best_task.pt" in text and "新版可选机制" in text
    assert "本轮保存 checkpoint_best_task.pt" not in text


def test_content_module_reexports_discussion_units_without_changing_legacy_build_input():
    from scripts.tho_group_meeting_ppt import build, content, discussion_content

    assert content.DISCUSSION_UNITS is discussion_content.DISCUSSION_UNITS
    assert build.MAIN_SLIDES is content.MAIN_SLIDES
    assert build.MAIN_SLIDES is not content.DISCUSSION_UNITS


def test_discussion_sources_are_existing_repository_files():
    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    for unit in DISCUSSION_UNITS:
        for source in unit.sources:
            base_path = source.split("#", 1)[0]
            assert not Path(base_path).is_absolute(), (unit.key, source)
            assert (REPO_ROOT / base_path).is_file(), (unit.key, source)


def test_discussion_content_avoids_deck_and_status_artifacts():
    from dataclasses import fields

    from scripts.tho_group_meeting_ppt.content import DISCUSSION_UNITS

    forbidden = ("工作进度", "实验编号汇总", "复现实验命令", "现场问答", "固定页数")
    text = "\n".join(
        str(value)
        for unit in DISCUSSION_UNITS
        for value in (getattr(unit, field.name) for field in fields(unit))
    )
    assert not any(phrase in text for phrase in forbidden)
    assert all(not hasattr(unit, "slide_number") for unit in DISCUSSION_UNITS)


def test_backup_deck_covers_required_materials():
    from scripts.tho_group_meeting_ppt.content import BACKUP_SLIDES

    keys = {slide.key for slide in BACKUP_SLIDES}
    assert {
        "balanced_cases",
        "model_details",
        "training_details",
        "metric_formulas",
        "full_stratified_results",
        "subject_results",
        "harmonic_subgroups",
        "harmonic_metrics",
        "data_provenance",
    } <= keys
    assert "reproduction_commands" not in keys
    assert "qa_notes" not in keys


def test_chart_data_matches_canonical_report_values():
    from scripts.tho_group_meeting_ppt.charts import load_chart_data

    data = load_chart_data(REPO_ROOT)

    assert data.overall.loc["g3_c_wide_8p0", "robust_rr"] == pytest.approx(0.712186)
    assert data.overall.loc["g3_c_bandenergy", "count_bpm"] == pytest.approx(0.684800)
    assert data.harmonic.loc["g3_c_bandenergy", "correction_rate"] == pytest.approx(0.9720, abs=5e-5)
    assert data.harmonic.loc["g3_c_wide_8p0", "lag4_corr"] == pytest.approx(0.847504)
    assert data.harmonic_coverage["positive_union_windows"] == 452


def _shape_text(slide, name: str) -> str:
    return next(shape.text for shape in slide.shapes if shape.name == name)


def test_generated_main_deck_has_25_slides_black_body_and_editable_diagrams(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.build import build_presentation
    from scripts.tho_group_meeting_ppt.theme import BODY_BLACK

    output = tmp_path / "main.pptx"
    build_presentation(template=TEMPLATE, output=output, include_backup=False)
    prs = Presentation(output)

    assert len(prs.slides) == 25
    assert _shape_text(prs.slides[10], "页面标题") == "独立测试集结果支持宽频 STFT 作为当前基准"
    assert all(slide.slide_layout.name != "空白" for slide in list(prs.slides)[1:])
    assert sum(1 for shape in prs.slides[6].shapes if shape.name.startswith("流程节点")) >= 8
    assert sum(1 for slide in prs.slides for shape in slide.shapes if shape.has_table) >= 3
    assert sum(1 for shape in prs.slides[10].shapes if shape.shape_type == 13) == 0
    assert sum(1 for shape in prs.slides[20].shapes if shape.has_table) == 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.name.startswith(("正文", "核心结论")) or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    assert run.font.color.rgb == BODY_BLACK


def test_backup_contains_balanced_cases_formulas_subject_table_and_commands(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.build import build_presentation

    output = tmp_path / "full.pptx"
    build_presentation(template=TEMPLATE, output=output, include_backup=True)
    prs = Presentation(output)
    text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )

    assert len(prs.slides) >= 36
    assert "dataset_row_id=640" in text
    assert "dataset_row_id=873" in text
    assert "dataset_row_id=1353" in text
    assert "dataset_row_id=3584" in text
    assert "8 名测试受试者" in text
    assert "8 名测试受试者（1/2）" in text
    assert "8 名测试受试者（2/2）" in text
    assert "复现实验命令与证据路径" not in text
    assert "现场问答提示" not in text
    assert "eval_topk_checkpoints.py" not in text


def test_final_deck_is_complete_editable_and_powerpoint_compatible():
    from zipfile import ZipFile

    from scripts.tho_group_meeting_ppt.detail_slides import SLIDE_GROUPS, UNIT_BY_KEY
    from scripts.tho_group_meeting_ppt.discussion_content import DISCUSSION_UNITS
    from scripts.tho_group_meeting_ppt.theme import BODY_BLACK, is_three_line_table

    expected_assets = {
        "signal_overview.png", "softz_mapping.png", "token_geometry.png",
        "stft_branch_shapes.png", "bandenergy_response.png", "loss_schedule.png",
        "metric_robust_rr.png", "metric_cycle_count.png", "metric_lag_corr.png",
        "metric_relative_envelope.png", "metric_local_rr.png", "overall_delta.png",
        "seed_subject_stability.png", "strata_tradeoffs.png",
        "stft_resolution_comparison.png", "case_row_640.png", "case_row_873.png",
        "case_row_1353.png", "case_row_3584.png",
    }
    forbidden = (
        "论文分享", "2021/12/21", "汇报人：xxx", "工作进度", "实验编号汇总",
        "复现实验命令", "现场问答", "更优",
    )
    body_exclusions = ("页面标题", "页码", "来源注释", "章节标签")

    assert hashlib.sha256(TEMPLATE.read_bytes()).hexdigest() == (
        "ce71b290bfd48cc74d7aef7f9586de05842c03349e933aa8f5a02572f2a72e10"
    )
    assert FINAL_DECK.exists()
    prs = Presentation(FINAL_DECK)
    assert len(DISCUSSION_UNITS) == 48
    assert len(prs.slides) == 1 + 11 + len(SLIDE_GROUPS) == 65

    all_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "THO research v2｜从整晚 BCG 到呼吸重建" in all_text
    assert "全流程研究方法细节讨论｜证据、边界与下一轮决策" in all_text
    assert not any(phrase in all_text for phrase in forbidden)
    with ZipFile(FINAL_DECK) as archive:
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels"))
        )
    assert not any(phrase in package_text for phrase in forbidden)

    sections = list(dict.fromkeys(spec.section for spec in SLIDE_GROUPS))
    assert len(sections) == 11
    section_pages: set[int] = set()
    page_index = 1
    current_section = None
    consumed_units = []
    for section_number, spec in enumerate(SLIDE_GROUPS, start=1):
        if spec.section != current_section:
            current_section = spec.section
            expected_number = sections.index(current_section) + 1
            section_pages.add(page_index)
            section_title = next(
                shape.text for shape in prs.slides[page_index].shapes
                if shape.name == "页面标题"
            )
            assert section_title == f"{expected_number:02d}｜{current_section}"
            page_index += 1
        slide = prs.slides[page_index]
        assert next(shape.text for shape in slide.shapes if shape.name == "页面标题") == spec.title
        slide_text = "\n".join(
            shape.text for shape in slide.shapes if shape.has_text_frame
        )
        keys = spec.unit_keys or spec.context_keys
        for key in keys:
            assert UNIT_BY_KEY[key].question in slide_text, (spec.key, key)
        consumed_units.extend(spec.unit_keys)
        page_index += 1
    assert page_index == len(prs.slides)
    assert consumed_units == list(dict.fromkeys(consumed_units))
    assert set(consumed_units) == {unit.key for unit in DISCUSSION_UNITS}

    actual_assets = set()
    image_count = 0
    native_formula_or_flow = 0
    tables = []
    case_pages = []
    for index, slide in enumerate(prs.slides):
        title = next(
            (shape.text for shape in slide.shapes if shape.name == "页面标题"), ""
        )
        if "dataset_row_id=" in title:
            case_pages.append(title)
        if index > 0 and index not in section_pages:
            assert title
            substantive = [
                shape for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
                and shape.name not in body_exclusions
            ]
            assert len(substantive) >= 2, (index + 1, title)
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= prs.slide_width
            assert shape.top + shape.height <= prs.slide_height
            if shape.shape_type == 13:
                image_count += 1
                actual_assets.add(shape.name.removeprefix("真实证据图："))
                assert shape.width >= Inches(4.5) and shape.height >= Inches(2.4)
                assert shape.width < prs.slide_width and shape.height < prs.slide_height
            if shape.name.startswith(("公式", "流程", "方法", "架构", "证据边界", "讨论框")):
                native_formula_or_flow += 1
            if shape.has_table:
                tables.append(shape.table)
                assert is_three_line_table(shape.table)
                for row in shape.table.rows:
                    for cell in row.cells:
                        tc_pr = cell._tc.get_or_add_tcPr()
                        assert all(
                            tc_pr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{edge}") is None
                            for edge in ("lnL", "lnR", "lnTlToBr", "lnBlToTr")
                        )
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    assert run.font.size.pt >= 14
            if shape.has_text_frame and shape.name not in body_exclusions and not shape.has_table:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            assert run.font.color.rgb == BODY_BLACK, (index + 1, shape.name, run.text)
                            assert run.font.size.pt >= 18, (index + 1, shape.name, run.text)

    assert image_count == len(expected_assets) == 19
    assert actual_assets == expected_assets
    assert native_formula_or_flow >= 35
    assert len(tables) == 5
    assert case_pages == [
        "dataset_row_id=640｜四模型均纠正",
        "dataset_row_id=873｜四模型均失败",
        "dataset_row_id=1353｜模型间分歧",
        "dataset_row_id=3584｜阈值边界",
    ]
    canonical = [
        table for table in tables
        if table._graphic_frame.name.startswith("四方案核心指标")
    ]
    assert len(canonical) == 2 and all(len(table.rows) == 5 for table in canonical)
    canonical_text = "\n".join(
        cell.text for table in canonical for row in table.rows for cell in row.cells
    )
    for scheme in ("纯时序", "F0 STFT", "wide STFT", "bandenergy"):
        assert canonical_text.count(scheme) == 2
    for metric in (
        "robust RR MAE bpm", "count bpm MAE", "lag4 corr（无量纲）",
        "relative envelope corr", "local RR MAE bpm",
    ):
        assert metric in canonical_text

    gaps = [
        shape for slide in prs.slides for shape in slide.shapes
        if shape.name.startswith("证据缺口")
    ]
    assert gaps
    assert all("所需字段：" in gap.text and "建议入口：" in gap.text for gap in gaps)
    assert all(_estimated_text_height(gap) <= gap.height / Inches(1) for gap in gaps)


def test_cli_can_run_from_repository_root():
    result = subprocess.run(
        [sys.executable, "scripts/build_tho_group_meeting_ppt.py", "--charts-only"],
        cwd=REPO_ROOT,
        capture_output=True,
        text=True,
        check=False,
    )
    assert result.returncode == 0, result.stderr


def test_discussion_slide_specs_consume_every_unit_once_and_cover_sections():
    from scripts.tho_group_meeting_ppt.detail_slides import SLIDE_GROUPS
    from scripts.tho_group_meeting_ppt.discussion_content import DISCUSSION_UNITS

    consumed = [key for spec in SLIDE_GROUPS for key in spec.unit_keys]
    assert len(SLIDE_GROUPS) >= 43
    assert consumed == list(dict.fromkeys(consumed))
    assert set(consumed) == {unit.key for unit in DISCUSSION_UNITS}
    assert {spec.section for spec in SLIDE_GROUPS} == EXPECTED_DISCUSSION_SECTIONS
    assert sum(len(spec.unit_keys) == 2 for spec in SLIDE_GROUPS) >= 3
    assert all(not hasattr(spec, "max_chars") for spec in SLIDE_GROUPS)


def test_discussion_deck_is_editable_readable_and_uses_three_line_tables(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.build import build_discussion_presentation
    from scripts.tho_group_meeting_ppt.theme import BODY_BLACK, is_three_line_table

    output = tmp_path / "discussion.pptx"
    template_hash = hashlib.sha256(TEMPLATE.read_bytes()).hexdigest()
    build_discussion_presentation(template=TEMPLATE, output=output, repo_root=REPO_ROOT)
    prs = Presentation(output)
    all_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
    )

    assert len(prs.slides) >= 55  # 标题 + 11 章导航 + 至少 43 张问题页
    for phrase in (
        "样本如何从整晚数据形成", "soft-z", "PatchMixer", "STFT",
            "L_signed_corr", "稳健 RR", "dataset_row_id=640", "最小判别性实验",
    ):
        assert phrase in all_text
    assert "20 秒局部 RR 窗" in all_text
    assert "40 秒窗 / 10 秒步长" in all_text
    assert "256 / 128 / 140" in all_text
    assert "本窗未触发压缩" in all_text
    for phrase in ("复现实验命令", "现场问答", "工作进度", "实验编号汇总", "泛化 Q&A", "更优"):
        assert phrase not in all_text

    picture_sizes = []
    native_formula_or_flow = 0
    tables = 0
    for index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= prs.slide_width
            assert shape.top + shape.height <= prs.slide_height
            if shape.shape_type == 13:
                picture_sizes.append((shape.width / Inches(1), shape.height / Inches(1)))
            if shape.name.startswith(("公式", "流程", "方法", "架构", "证据边界", "讨论框")):
                native_formula_or_flow += 1
            if shape.has_table:
                tables += 1
                assert is_three_line_table(shape.table)
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                assert run.font.size.pt >= 14
            if shape.has_text_frame and shape.name.startswith(("正文", "方法", "证据", "限制", "讨论")):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            assert run.font.color.rgb == BODY_BLACK
                            assert run.font.size.pt >= 18
        if index >= 12:  # 非标题、非章节页
            assert any(shape.name == "页面标题" and shape.text.strip() for shape in slide.shapes)
            assert sum(bool(shape.has_text_frame and shape.text.strip()) for shape in slide.shapes) >= 2
            assert not (len(slide.shapes) == 1 and slide.shapes[0].shape_type == 13)
    assert tables >= 3
    assert native_formula_or_flow >= 35
    assert picture_sizes and all(w >= 4.5 and h >= 2.4 for w, h in picture_sizes)
    assert hashlib.sha256(TEMPLATE.read_bytes()).hexdigest() == template_hash


def test_discussion_evidence_gap_is_native_and_actionable(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.build import build_discussion_presentation

    output = tmp_path / "discussion_gap.pptx"
    build_discussion_presentation(template=TEMPLATE, output=output, repo_root=REPO_ROOT)
    prs = Presentation(output)
    gaps = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.name.startswith("证据缺口")
    ]
    assert gaps
    assert all("所需字段：" in text and "建议入口：" in text for text in gaps)


def test_cli_discussion_routes_are_mutually_exclusive_and_default_to_discussion(tmp_path: Path):
    script = REPO_ROOT / "scripts/build_tho_group_meeting_ppt.py"
    result = subprocess.run(
        [sys.executable, str(script), "--evidence-only", "--assets-only"],
        cwd=REPO_ROOT, capture_output=True, text=True, check=False,
    )
    assert result.returncode != 0
    output = tmp_path / "default_discussion.pptx"
    result = subprocess.run(
        [sys.executable, str(script), "--output", str(output)],
        cwd=REPO_ROOT, capture_output=True, text=True, check=False,
    )
    assert result.returncode == 0, result.stderr
    assert output.is_file() and len(Presentation(output).slides) >= 55


def test_cli_default_output_is_the_documented_formal_delivery(monkeypatch):
    import importlib

    monkeypatch.syspath_prepend(str(REPO_ROOT / "scripts"))
    cli = importlib.import_module("build_tho_group_meeting_ppt")
    assert cli.build_parser().parse_args([]).output == FINAL_DECK

    readme = (REPO_ROOT / "docs/stage_reports/README.md").read_text(encoding="utf-8")
    assert "./.venv/bin/python scripts/build_tho_group_meeting_ppt.py" in readme
    assert str(FINAL_DECK.relative_to(REPO_ROOT)) in readme


@pytest.fixture(scope="module")
def deterministic_discussion_decks(tmp_path_factory):
    from scripts.tho_group_meeting_ppt.build import build_discussion_presentation

    root = tmp_path_factory.mktemp("deterministic_discussion")
    outputs = (root / "first.pptx", root / "second.pptx")
    for output in outputs:
        build_discussion_presentation(template=TEMPLATE, output=output, repo_root=REPO_ROOT)
    return outputs


def test_discussion_pptx_is_byte_deterministic_and_matches_formal_file(
    deterministic_discussion_decks,
):
    from zipfile import ZIP_DEFLATED, ZipFile

    first, second = deterministic_discussion_decks
    hashes = [hashlib.sha256(path.read_bytes()).hexdigest() for path in (first, second, FINAL_DECK)]
    assert len(set(hashes)) == 1
    for path in (first, second, FINAL_DECK):
        with ZipFile(path) as archive:
            infos = archive.infolist()
            assert [info.filename for info in infos] == sorted(info.filename for info in infos)
            assert all(info.date_time == (1980, 1, 1, 0, 0, 0) for info in infos)
            assert all(info.create_system == 0 for info in infos)
            assert all(info.external_attr == 0x20 for info in infos)
            assert all(info.compress_type == ZIP_DEFLATED for info in infos)


def test_discussion_core_metadata_is_project_owned_and_reproducible(
    deterministic_discussion_decks,
):
    from xml.etree import ElementTree
    from zipfile import ZipFile

    namespaces = {
        "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
        "dc": "http://purl.org/dc/elements/1.1/",
        "dcterms": "http://purl.org/dc/terms/",
    }
    expected = {
        "dc:creator": "THO research v2",
        "cp:lastModifiedBy": "THO research v2",
        "dcterms:created": "2026-07-08T00:00:00Z",
        "dcterms:modified": "2026-07-08T00:00:00Z",
    }
    for path in (*deterministic_discussion_decks, FINAL_DECK):
        with ZipFile(path) as archive:
            core = archive.read("docProps/core.xml")
        text = core.decode("utf-8")
        assert not any(value in text for value in ("miao mengya", "MISTGUN", "2019", "2022"))
        root = ElementTree.fromstring(core)
        for field, value in expected.items():
            assert root.find(field, namespaces).text == value


@pytest.mark.parametrize(
    ("first", "second"),
    [
        (first, second)
        for index, first in enumerate((
            "--evidence-only", "--assets-only", "--discussion-deck",
            "--legacy-summary", "--charts-only",
        ))
        for second in (
            "--evidence-only", "--assets-only", "--discussion-deck",
            "--legacy-summary", "--charts-only",
        )[index + 1:]
    ],
)
def test_all_cli_modes_are_pairwise_mutually_exclusive(first: str, second: str):
    script = REPO_ROOT / "scripts/build_tho_group_meeting_ppt.py"
    result = subprocess.run(
        [sys.executable, str(script), first, second],
        cwd=REPO_ROOT, capture_output=True, text=True, check=False,
    )
    assert result.returncode != 0


def test_discussion_specs_have_explicit_content_plans_and_no_silent_truncation():
    import inspect

    from scripts.tho_group_meeting_ppt import detail_slides

    assert all(spec.display_fields for spec in detail_slides.SLIDE_GROUPS)
    assert all(spec.panel_plan for spec in detail_slides.SLIDE_GROUPS)
    allowed = {
        "question", "method_steps", "parameters", "rationale", "evidence",
        "limits", "discussion_prompt",
    }
    assert all(set(spec.display_fields) <= allowed for spec in detail_slides.SLIDE_GROUPS)
    for spec in detail_slides.SLIDE_GROUPS:
        units = [
            detail_slides.UNIT_BY_KEY[key]
            for key in (spec.unit_keys or spec.context_keys)
        ]
        for panel in spec.panel_plan:
            if panel.field == "discussion_prompt":
                values = [unit.discussion_prompt for unit in units if unit.discussion_prompt]
            else:
                values = [value for unit in units for value in getattr(unit, panel.field)]
            assert values, (spec.key, panel.field)
    source = inspect.getsource(detail_slides)
    for forbidden in ("limit=3", "limit=4", "steps[:4]", "panels[:3]", "values[:limit]"):
        assert forbidden not in source


def test_canonical_four_model_table_reads_manifest_backed_five_metrics():
    from scripts.tho_group_meeting_ppt.detail_slides import load_canonical_overall_metrics

    rows = load_canonical_overall_metrics(REPO_ROOT)
    assert [row["label"] for row in rows] == [
        "g0_time_only", "g0_f0_native_stft_pre_mixer", "g3_c_wide_8p0", "g3_c_bandenergy",
    ]
    assert all(set(row) == {
        "label", "robust_rr_mae_bpm", "count_bpm_mae", "lag4_corr",
        "relative_envelope_corr", "local_rr_mae_bpm",
    } for row in rows)
    assert rows[0]["relative_envelope_corr"] == pytest.approx(0.596490)
    assert rows[2]["robust_rr_mae_bpm"] == pytest.approx(0.712186)
    assert rows[3]["count_bpm_mae"] == pytest.approx(0.684800)


@pytest.fixture(scope="module")
def revised_discussion_deck(tmp_path_factory):
    from scripts.tho_group_meeting_ppt.build import build_discussion_presentation

    output = tmp_path_factory.mktemp("revised_discussion") / "discussion.pptx"
    build_discussion_presentation(template=TEMPLATE, output=output, repo_root=REPO_ROOT)
    return Presentation(output)


def _bbox_overlap_ratio(a, b) -> float:
    left = max(a.left, b.left)
    top = max(a.top, b.top)
    right = min(a.left + a.width, b.left + b.width)
    bottom = min(a.top + a.height, b.top + b.height)
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    return intersection / min(a.width * a.height, b.width * b.height)


def _estimated_text_height(shape) -> float:
    from PIL import ImageFont

    runs = [
        run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs
        if run.text.strip()
    ]
    if not runs:
        return 0.0
    size_pt = max((run.font.size.pt if run.font.size else 18.0) for run in runs)
    font = ImageFont.truetype("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", round(size_pt * 96 / 72))
    width_px = max(1.0, shape.width / Inches(1) * 96 - 24)
    height_inches = (shape.text_frame.margin_top + shape.text_frame.margin_bottom) / Inches(1)
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        if not text:
            continue
        current = ""
        paragraph_lines = 1
        for character in text:
            candidate = current + character
            if current and font.getlength(candidate) > width_px:
                paragraph_lines += 1
                current = character
            else:
                current = candidate
        paragraph_size = max(
            (run.font.size.pt if run.font.size else size_pt)
            for run in paragraph.runs
        )
        height_inches += paragraph_lines * paragraph_size * 1.32 / 72.0
        if paragraph.space_before:
            height_inches += paragraph.space_before.pt / 72.0
        if paragraph.space_after:
            height_inches += paragraph.space_after.pt / 72.0
    return height_inches


def test_critical_discussion_pages_have_no_text_overlap_or_estimated_overflow(revised_discussion_deck):
    title_fragments = (
        "样本如何从整晚数据形成", "soft-z", "目标、整窗筛选", "目标监督代码证据",
        "目标监督边界", "PatchMixer",
        "WeakSyncLoss", "checkpoint", "四方案统一比较（一）", "四方案统一比较（二）",
        "local RR", "dataset_row_id=640",
        "dataset_row_id=873", "dataset_row_id=1353", "dataset_row_id=3584",
        "最小判别性实验",
    )
    checked = set()
    for slide in revised_discussion_deck.slides:
        title = next((shape.text for shape in slide.shapes if shape.name == "页面标题"), "")
        if not any(fragment in title for fragment in title_fragments):
            continue
        checked.add(next(fragment for fragment in title_fragments if fragment in title))
        content = [
            shape for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
            and shape.name.startswith((
                "正文", "方法", "证据", "限制", "讨论", "副标题",
                "流程节点文字", "完整复核",
            ))
        ]
        for index, first in enumerate(content):
            assert _estimated_text_height(first) <= first.height / Inches(1), (title, first.name, first.text)
            for second in content[index + 1:]:
                assert _bbox_overlap_ratio(first, second) < 0.15, (title, first.name, second.name)
    assert checked == set(title_fragments)


def test_discussion_deck_embeds_all_real_assets_and_complete_canonical_tables(revised_discussion_deck):
    expected = {
        "signal_overview.png", "softz_mapping.png", "token_geometry.png", "stft_branch_shapes.png",
        "bandenergy_response.png", "loss_schedule.png", "metric_robust_rr.png", "metric_cycle_count.png",
        "metric_lag_corr.png", "metric_relative_envelope.png", "metric_local_rr.png", "overall_delta.png",
        "seed_subject_stability.png", "strata_tradeoffs.png", "stft_resolution_comparison.png",
        "case_row_640.png", "case_row_873.png", "case_row_1353.png", "case_row_3584.png",
    }
    actual = {
        shape.name.removeprefix("真实证据图：")
        for slide in revised_discussion_deck.slides
        for shape in slide.shapes
        if shape.shape_type == 13
    }
    assert actual == expected
    tables = [
        shape.table for slide in revised_discussion_deck.slides for shape in slide.shapes
        if shape.has_table and shape.name.startswith("四方案核心指标")
    ]
    assert len(tables) == 2
    assert all(len(table.rows) == 5 for table in tables)
    table_text = "\n".join(cell.text for table in tables for row in table.rows for cell in row.cells)
    for phrase in ("纯时序", "F0 STFT", "wide STFT", "bandenergy", "relative envelope", "local RR"):
        assert phrase in table_text


def _copy_discussion_assets(tmp_path: Path) -> Path:
    import shutil

    source = REPO_ROOT / "docs/stage_reports/20260708/generated_assets/discussion"
    target = tmp_path / "discussion"
    target.mkdir()
    for path in source.iterdir():
        if path.suffix in {".png", ".json"}:
            shutil.copy2(path, target / path.name)
    return target


def test_discussion_asset_resolver_validates_all_19_manifested_assets(tmp_path: Path):
    from scripts.tho_group_meeting_ppt.detail_slides import resolve_discussion_assets

    assets = resolve_discussion_assets(REPO_ROOT, asset_dir=_copy_discussion_assets(tmp_path))
    assert len(assets) == 19


@pytest.mark.parametrize("tamper", ["asset", "generator", "source"])
def test_discussion_asset_resolver_rejects_tampered_evidence(tmp_path: Path, tamper: str):
    from scripts.tho_group_meeting_ppt.detail_slides import resolve_discussion_assets

    asset_dir = _copy_discussion_assets(tmp_path)
    manifest_path = asset_dir / "signal_assets_manifest.json"
    if tamper == "asset":
        (asset_dir / "signal_overview.png").write_bytes(b"tampered")
    else:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        key = "signal_figure_code" if tamper == "generator" else "signal_metadata_json"
        source = Path(manifest["evidence"][key]["path"])
        source = source if source.is_absolute() else REPO_ROOT / source
        tampered_source = tmp_path / source.name
        tampered_source.write_bytes(source.read_bytes() + b"\ntampered")
        manifest["evidence"][key]["path"] = str(tampered_source)
        manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
    with pytest.raises(ValueError, match=r"signal_assets_manifest\.json.*--assets-only"):
        resolve_discussion_assets(REPO_ROOT, asset_dir=asset_dir)


def test_case_titles_respect_logo_safe_area_and_local_rr_has_single_evidence_boundary(revised_discussion_deck):
    local_rr_slides = []
    for slide in revised_discussion_deck.slides:
        title_shape = next((shape for shape in slide.shapes if shape.name == "页面标题"), None)
        if title_shape is None:
            continue
        if "dataset_row_id=" in title_shape.text:
            assert title_shape.text.startswith("dataset_row_id=")
            assert len(title_shape.text) <= 31
            assert title_shape.left + title_shape.width <= Inches(11.90)
        if title_shape.text.startswith("local RR：正式"):
            local_rr_slides.append(slide)
    assert len(local_rr_slides) == 1
    assert sum(shape.name == "证据边界" for shape in local_rr_slides[0].shapes) == 1


@pytest.mark.skipif(
    os.environ.get("RUN_LIBREOFFICE_INTEGRATION") != "1",
    reason="设置 RUN_LIBREOFFICE_INTEGRATION=1 后显式运行 LibreOffice 集成测试",
)
def test_key_pages_survive_real_libreoffice_pdf_and_png_render(tmp_path: Path):
    from PIL import Image, ImageStat

    from scripts.tho_group_meeting_ppt.build import build_discussion_presentation

    deck = tmp_path / "discussion.pptx"
    build_discussion_presentation(template=TEMPLATE, output=deck, repo_root=REPO_ROOT)
    rendered = tmp_path / "rendered"
    rendered.mkdir()
    lo_home = tmp_path / "lohome"
    runtime = tmp_path / "runtime"
    profile = tmp_path / "lo_profile"
    lo_home.mkdir()
    runtime.mkdir(mode=0o700)
    profile.mkdir()
    lo_env = os.environ.copy()
    lo_env.update({
        "HOME": str(lo_home),
        "XDG_RUNTIME_DIR": str(runtime),
        "TMPDIR": str(tmp_path),
        "SAL_USE_VCLPLUGIN": "svp",
    })
    conversion = subprocess.run(
        [
            "libreoffice", f"-env:UserInstallation={profile.as_uri()}", "--headless",
            "--nologo", "--nodefault", "--nofirststartwizard", "--norestore",
            "--convert-to", "pdf", "--outdir", str(rendered), str(deck),
        ],
        capture_output=True, text=True, check=False, env=lo_env,
    )
    assert conversion.returncode == 0, (
        f"LibreOffice 转换失败: returncode={conversion.returncode}\n"
        f"stdout={conversion.stdout}\nstderr={conversion.stderr}\n"
        f"HOME={lo_env['HOME']}\nXDG_RUNTIME_DIR={lo_env['XDG_RUNTIME_DIR']}\n"
        f"profile={profile.as_uri()}"
    )
    pdf = rendered / "discussion.pdf"
    assert pdf.is_file()
    extracted = subprocess.run(
        ["pdftotext", "-f", "14", "-l", "56", "-layout", str(pdf), "-"],
        capture_output=True, text=True, check=False,
    )
    assert extracted.returncode == 0, extracted.stderr
    for phrase in (
        "本窗 0 samples changed", "checkpoint_top1/2/3", "固定类别与 row",
        "同一 seed 四模型", "还需哪些反例类别", "gate 风",
    ):
        assert phrase in extracted.stdout
    for page in (14, 29, 53, 54, 55, 56):
        prefix = rendered / f"page_{page}"
        raster = subprocess.run(
            ["pdftoppm", "-f", str(page), "-l", str(page), "-singlefile", "-png", "-r", "96", str(pdf), str(prefix)],
            capture_output=True, text=True, check=False,
        )
        assert raster.returncode == 0, raster.stderr
        image = Image.open(prefix.with_suffix(".png")).convert("RGB")
        assert ImageStat.Stat(image).mean[0] > 200
        assert ImageStat.Stat(image).stddev[0] > 15


@pytest.fixture(scope="module")
def model_asset_builds(tmp_path_factory):
    from scripts.tho_group_meeting_ppt.model_figures import build_model_assets

    first_dir = tmp_path_factory.mktemp("model_assets_first")
    second_dir = tmp_path_factory.mktemp("model_assets_second")
    first = build_model_assets(REPO_ROOT, first_dir)
    second = build_model_assets(REPO_ROOT, second_dir)
    return first, second


def test_token_geometry_comes_from_formal_model_and_actual_patch_mixer(model_asset_builds):
    import torch

    from resp_train.models.timeseries import PatchMixer1D

    (assets, metadata), (_, repeated_metadata) = model_asset_builds
    model = PatchMixer1D(
        base_channels=metadata["token_shape"][1],
        patch_len=metadata["patch_len"],
        patch_stride=metadata["patch_stride"],
        mixer_layers=metadata["mixer_layers"],
    )
    tokens, original_length = model.tokenize_input(torch.zeros(2, 1, 18_000))

    assert set(assets) == {"token_geometry", "stft_branch_shapes", "loss_schedule"}
    assert metadata["patch_len"] == 256
    assert metadata["patch_stride"] == 128
    assert metadata["patch_count"] == model.token_count_for_length(18_000) == 140
    assert metadata["token_shape"] == ["B", 16, 140]
    assert list(tokens.shape) == [2, 16, 140]
    assert original_length == 18_000
    assert metadata["token_shape"] == repeated_metadata["token_shape"]
    assert "g_series_stft_input" in Path(metadata["sources"]["formal_config"]).parts


def test_model_detail_stft_shapes_match_centered_forward_and_formal_g_configs(model_asset_builds):
    from PIL import Image

    (assets, metadata), (repeated_assets, repeated_metadata) = model_asset_builds
    manifest = json.loads((next(iter(assets.values())).parent / "model_assets_manifest.json").read_text())
    generator = manifest["sources"]["model_figure_code"]
    generator_path = REPO_ROOT / generator["path"]
    assert generator_path.name == "model_figures.py"
    assert generator["sha256"] == hashlib.sha256(generator_path.read_bytes()).hexdigest()
    branches = metadata["stft_branches"]

    assert set(branches) == {"f0", "wide", "bandenergy"}
    assert branches["f0"]["raw_stft_shape"] == ["B", 1501, 37]
    assert branches["f0"]["cropped_shape"] == ["B", 239, 37]
    assert branches["wide"]["raw_stft_shape"] == ["B", 1001, 73]
    assert branches["wide"]["cropped_shape"] == ["B", 160, 73]
    assert branches["bandenergy"]["encoder_input_shape"] == ["B", 5, 73]
    assert all(branch["encoder_output_shape"][1:] == [16, branch["raw_stft_shape"][-1]] for branch in branches.values())
    assert all(branch["aligned_token_shape"] == ["B", 16, 140] for branch in branches.values())
    assert all(branch["inject_position"] == "pre_mixer" for branch in branches.values())
    assert branches["f0"]["encoder_type"] == branches["wide"]["encoder_type"] == "conv2d"
    assert branches["bandenergy"]["encoder_type"] == "bandenergy"
    assert all(branch["zero_and_real_probe_shapes_match"] for branch in branches.values())
    assert metadata["stft_branches"] == repeated_metadata["stft_branches"]

    for key, path in assets.items():
        assert Image.open(path).size[0] >= 1600
        assert Image.open(path).size[1] >= 900
        assert metadata["layout"][key]["all_key_text_inside"] is True
        assert metadata["layout"][key]["text_overlap_count"] == 0
        assert hashlib.sha256(path.read_bytes()).hexdigest() == hashlib.sha256(
            repeated_assets[key].read_bytes()
        ).hexdigest()


def test_loss_schedule_is_evaluated_by_weak_sync_loss_from_formal_config(model_asset_builds):
    from omegaconf import OmegaConf

    from resp_train.losses.weak import WeakSyncLoss

    (_, metadata), _ = model_asset_builds
    formal_config = Path(metadata["sources"]["formal_config"])
    cfg = OmegaConf.load(formal_config)
    loss = WeakSyncLoss(cfg)
    actual = {"signed_corr": [], "signed_cosine": []}
    for epoch in range(1, 11):
        loss.set_epoch(epoch)
        weights = loss.current_weights()
        actual["signed_corr"].append(weights["signed_corr"])
        actual["signed_cosine"].append(weights["signed_cosine"])

    schedule = metadata["loss_schedule"]
    assert schedule["signed_corr"] == pytest.approx(actual["signed_corr"])
    assert schedule["signed_cosine"] == pytest.approx(actual["signed_cosine"])
    assert schedule["signed_corr"][:6] == pytest.approx([0.6, 0.52, 0.44, 0.36, 0.28, 0.2])
    assert schedule["signed_cosine"][:6] == pytest.approx([0.1, 0.08, 0.06, 0.04, 0.02, 0.0])
    assert schedule["fixed_weights"]["envelope"] == pytest.approx(1.0)
    assert schedule["fixed_weights"]["spectrum"] == pytest.approx(0.2)
    assert "stft_dist" in schedule["disabled_terms"]
    assert Path(metadata["sources"]["loss_code"]).samefile(REPO_ROOT / "resp_train/losses/weak.py")


def test_loss_schedule_classifies_every_weighted_weak_sync_term(model_asset_builds):
    import ast
    import inspect
    import textwrap

    from resp_train.losses.weak import WeakSyncLoss

    (_, metadata), _ = model_asset_builds
    schedule = metadata["loss_schedule"]
    expected = {
        "envelope", "spectrum", "smooth", "high_freq", "relative_envelope",
        "phase_alignment", "band_waveform", "curvature", "rhythm",
        "signed_cosine", "signed_corr", "signed_rms_envelope", "signed_mean",
        "si_sdr", "stft_dist", "stft_band_energy", "stft_peak_anchor",
        "fb_aux", "fb_consistency",
    }

    assert set(schedule["all_weighted_terms"]) == expected
    terms_in_weighted_total = set()
    for method in (
        WeakSyncLoss._weighted_base_total,
        WeakSyncLoss._weighted_stft_total,
        WeakSyncLoss._weighted_fb_total,
    ):
        tree = ast.parse(textwrap.dedent(inspect.getsource(method)))
        terms_in_weighted_total.update(
            node.slice.value
            for node in ast.walk(tree)
            if isinstance(node, ast.Subscript)
            and isinstance(node.value, ast.Name)
            and node.value.id == "components"
            and isinstance(node.slice, ast.Constant)
            and isinstance(node.slice.value, str)
        )
    assert terms_in_weighted_total == expected
    assert set(schedule["fixed_weights"]) | set(schedule["scheduled_terms"]) | set(
        schedule["disabled_terms"]
    ) == expected
    assert not (
        set(schedule["fixed_weights"]) & set(schedule["scheduled_terms"])
        or set(schedule["fixed_weights"]) & set(schedule["disabled_terms"])
        or set(schedule["scheduled_terms"]) & set(schedule["disabled_terms"])
    )
    assert "curvature" in schedule["disabled_terms"]


def test_loss_schedule_rejects_any_formal_g_config_drift(tmp_path: Path):
    from dataclasses import replace

    from omegaconf import OmegaConf

    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog
    from scripts.tho_group_meeting_ppt.model_figures import validate_formal_loss_configs

    catalog = build_evidence_catalog(REPO_ROOT)
    label = "g3_c_wide_8p0"
    drift_path = tmp_path / "drift.yaml"
    cfg = OmegaConf.load(catalog.run_configs[label].path)
    OmegaConf.update(cfg, "loss.curvature_weight", 0.3, merge=False)
    OmegaConf.save(cfg, drift_path)
    drift_configs = dict(catalog.run_configs)
    drift_configs[label] = replace(drift_configs[label], path=drift_path)
    drift_catalog = replace(catalog, run_configs=drift_configs)

    with pytest.raises(
        ValueError,
        match=r"g3_c_wide_8p0.*loss\.curvature_weight.*drift\.yaml",
    ):
        validate_formal_loss_configs(drift_catalog)


@pytest.fixture(scope="module")
def metric_asset_build(tmp_path_factory):
    from scripts.tho_group_meeting_ppt.metric_figures import build_metric_assets

    output = tmp_path_factory.mktemp("metric_assets")
    return build_metric_assets(REPO_ROOT, output)


def test_metric_assets_use_row_8025_and_current_repository_algorithms(metric_asset_build):
    import numpy as np

    from resp_train.metrics.signal import (
        bandpass_filter,
        best_lag_correlation_from_filtered,
        estimate_robust_peak_rate_bpm,
        lag_aligned_overlap,
        local_rr_metrics,
        relative_envelope_metrics,
        zero_crossing_counts,
    )
    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog

    assets, values = metric_asset_build
    catalog = build_evidence_catalog(REPO_ROOT)
    with np.load(catalog.general_signal_npz, allow_pickle=False) as data:
        pred = np.asarray(data["f0_prediction_full"], dtype=np.float64)
        target = np.asarray(data["target_respiration_full"], dtype=np.float64)
    pred_filtered = bandpass_filter(pred, fs=100.0, low_hz=0.05, high_hz=0.7, order=4)
    target_filtered = bandpass_filter(target, fs=100.0, low_hz=0.05, high_hz=0.7, order=4)
    pred_rr = estimate_robust_peak_rate_bpm(pred_filtered, fs=100.0, low_hz=0.05, high_hz=0.7)
    target_rr = estimate_robust_peak_rate_bpm(target_filtered, fs=100.0, low_hz=0.05, high_hz=0.7)
    pred_count = zero_crossing_counts(pred_filtered)
    target_count = zero_crossing_counts(target_filtered)
    lag = best_lag_correlation_from_filtered(
        pred_filtered, target_filtered, fs=100.0, max_lag_sec=4.0, low_hz=0.05
    )
    pred_overlap, target_overlap = lag_aligned_overlap(
        pred, target, lag_samples=round(lag["best_lag_sec"] * 100.0)
    )
    rel_env = relative_envelope_metrics(
        pred_overlap, target_overlap, fs=100.0, envelope_window_sec=2.0
    )
    local = local_rr_metrics(
        pred_filtered,
        target_filtered,
        fs=100.0,
        window_sec=40.0,
        step_sec=10.0,
        low_hz=0.05,
        high_hz=0.7,
    )

    assert set(assets) == {
        "metric_robust_rr",
        "metric_cycle_count",
        "metric_lag_corr",
        "metric_relative_envelope",
        "metric_local_rr",
    }
    assert values["dataset_row_id"] == 8025
    assert values["sample_split"] == "val"
    assert values["sample_model_label"] == "F0_native_stft_pre_mixer"
    assert values["robust_rr_abs_error"] == pytest.approx(abs(pred_rr - target_rr), abs=1e-6)
    assert values["cycle_count_abs_error"] == abs(pred_count["cycle"] - target_count["cycle"])
    assert values["cycle_count_abs_error"] >= 0
    assert values["best_lag_corr_4s"] == pytest.approx(lag["best_lag_corr"], abs=1e-9)
    assert -1.0 <= values["best_lag_corr_4s"] <= 1.0
    assert values["relative_envelope_mae_lag4s"] == pytest.approx(
        rel_env["relative_envelope_mae"], abs=1e-9
    )
    assert values["local_rr_mae"] == pytest.approx(local["local_rr_mae"], abs=1e-9)
    assert 0.0 <= values["local_rr_valid_frac"] <= 1.0
    assert values["parameters"]["local_rr_window_sec"] == 40.0
    assert values["parameters"]["local_rr_step_sec"] == 10.0
    assert values["parameters"]["lag_trace_points"] == 801
    assert values["parameters"]["lag_trace_integer_samples"] == [-400, 400]
    assert values["canonical_csv_alignment"] is False
    assert "计算示例" in values["evidence_scope"]
    assert "更优" not in json.dumps(values, ensure_ascii=False)


def test_metric_assets_are_large_uncropped_and_manifested(metric_asset_build):
    from PIL import Image, ImageStat

    assets, values = metric_asset_build
    output = next(iter(assets.values())).parent
    manifest = json.loads((output / "metric_assets_manifest.json").read_text(encoding="utf-8"))

    assert manifest["dataset_row_id"] == 8025
    assert manifest["evidence_scope"] == values["evidence_scope"]
    assert manifest["sources"]["signal_npz"]["path"].endswith("f0_visual_sample_signals.npz")
    assert manifest["sources"]["source_metrics_csv"]["path"].endswith("metrics.csv")
    assert manifest["sources"]["metric_code"]["path"] == "resp_train/metrics/signal.py"
    generator = REPO_ROOT / "scripts/tho_group_meeting_ppt/metric_figures.py"
    assert manifest["sources"]["metric_figure_code"]["path"] == (
        "scripts/tho_group_meeting_ppt/metric_figures.py"
    )
    assert manifest["sources"]["metric_figure_code"]["status"] == "present"
    assert manifest["sources"]["metric_figure_code"]["sha256"] == hashlib.sha256(
        generator.read_bytes()
    ).hexdigest()
    assert manifest["sources"]["robust_rr_demo_code"]["path"].endswith(
        "plot_rr_peak_band_metric_demo.py"
    )
    assert manifest["sources"]["signal_assets_manifest"]["path"].endswith(
        "signal_assets_manifest.json"
    )
    assert all(not Path(record["path"]).is_absolute() for record in manifest["sources"].values())
    assert manifest["parameters"]["local_rr_window_sec"] == 40.0
    assert manifest["parameters"]["local_rr_step_sec"] == 10.0
    assert set(manifest["assets"]) == set(assets)
    for key, path in assets.items():
        assert path.is_file() and path.stat().st_size > 20_000
        assert manifest["assets"][key]["sha256"] == hashlib.sha256(path.read_bytes()).hexdigest()
        assert values["layout"][key]["all_key_text_inside"] is True
        assert values["layout"][key]["text_overlap_count"] == 0
        with Image.open(path) as image:
            assert image.width >= 1600 and image.height >= 900
            assert ImageStat.Stat(image.convert("L")).stddev[0] > 10


@pytest.mark.parametrize(
    ("field", "bad_value"),
    (("model_label", "g3_c_wide_8p0"), ("split", "test"), ("dataset_row_id", 5058)),
)
def test_metric_asset_identity_rejects_tampered_general_signal_metadata(field, bad_value):
    import pandas as pd

    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog
    from scripts.tho_group_meeting_ppt.metric_figures import validate_metric_asset_identity

    catalog = build_evidence_catalog(REPO_ROOT)
    metadata_path = catalog.general_signal_npz.with_name("f0_visual_sample_metadata.json")
    metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
    metadata[field] = bad_value
    source_config = (REPO_ROOT / metadata["source_run_config"]).resolve()
    source_checkpoint = (REPO_ROOT / metadata["source_checkpoint"]).resolve()
    source_metrics = source_config.parent / "metrics.csv"
    frame = pd.read_csv(source_metrics)
    source_row = frame[frame["dataset_row_id"].astype(int) == 8025].iloc[0]

    with pytest.raises(ValueError, match=field):
        validate_metric_asset_identity(
            metadata,
            signal_npz=catalog.general_signal_npz,
            source_config=source_config,
            source_checkpoint=source_checkpoint,
            source_metrics_row=source_row,
        )


@pytest.mark.parametrize(
    ("tamper", "message"),
    (
        ("config", r"model\.stft_win.*config\.yaml"),
        ("csv", r"input_set.*metrics\.csv"),
        ("npz_hash", r"signal_npz\.sha256.*signal_assets_manifest\.json"),
        ("metadata_source", r"source_checkpoint.*metadata\.json"),
    ),
)
def test_metric_source_evidence_rejects_semantic_or_frozen_hash_drift(tmp_path: Path, tamper, message):
    import pandas as pd
    from omegaconf import OmegaConf

    from scripts.tho_group_meeting_ppt.evidence import build_evidence_catalog
    from scripts.tho_group_meeting_ppt.metric_figures import validate_metric_source_evidence

    catalog = build_evidence_catalog(REPO_ROOT)
    actual_metadata = catalog.general_signal_npz.with_name("f0_visual_sample_metadata.json")
    metadata = json.loads(actual_metadata.read_text(encoding="utf-8"))
    actual_config = (REPO_ROOT / metadata["source_run_config"]).resolve()
    actual_checkpoint = (REPO_ROOT / metadata["source_checkpoint"]).resolve()
    actual_csv = actual_config.parent / "metrics.csv"
    upstream = json.loads(
        (REPO_ROOT / "docs/stage_reports/20260708/generated_assets/discussion/signal_assets_manifest.json")
        .read_text(encoding="utf-8")
    )

    metadata_path = tmp_path / "metadata.json"
    config_path = tmp_path / "config.yaml"
    csv_path = tmp_path / "metrics.csv"
    upstream_path = tmp_path / "signal_assets_manifest.json"
    OmegaConf.save(OmegaConf.load(actual_config), config_path)
    metadata["source_run_config"] = str(config_path)
    frame = pd.read_csv(actual_csv)
    frame[frame["dataset_row_id"].astype(int) == 8025].to_csv(csv_path, index=False)

    if tamper == "config":
        cfg = OmegaConf.load(config_path)
        OmegaConf.update(cfg, "model.stft_win", 2000, merge=False)
        OmegaConf.save(cfg, config_path)
    elif tamper == "csv":
        changed = pd.read_csv(csv_path)
        changed.loc[0, "input_set"] = "other_input"
        changed.to_csv(csv_path, index=False)
    elif tamper == "metadata_source":
        metadata["source_checkpoint"] = "runs/other/checkpoint.pt"
    metadata_path.write_text(json.dumps(metadata, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    def record(path: Path) -> dict[str, object]:
        payload = path.read_bytes()
        return {
            "path": str(path),
            "status": "present",
            "sha256": hashlib.sha256(payload).hexdigest(),
            "size_bytes": len(payload),
        }

    upstream["evidence"]["signal_metadata_json"] = record(metadata_path)
    upstream["evidence"]["source_run_config"] = record(config_path)
    if tamper == "npz_hash":
        upstream["evidence"]["signal_npz"]["sha256"] = "0" * 64
    upstream_path.write_text(json.dumps(upstream, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    with pytest.raises(ValueError, match=message):
        validate_metric_source_evidence(
            repo_root=REPO_ROOT,
            signal_npz=catalog.general_signal_npz,
            signal_metadata_json=metadata_path,
            source_config=config_path,
            source_checkpoint=actual_checkpoint,
            source_metrics_csv=csv_path,
            signal_assets_manifest=upstream_path,
        )


def test_model_detail_full_registry_models_validate_native_injection_chain(model_asset_builds):
    import torch
    from omegaconf import OmegaConf

    from resp_train.models.registry import build_model

    (_, metadata), _ = model_asset_builds
    for key in ("f0", "wide", "bandenergy"):
        branch = metadata["stft_branches"][key]
        cfg = OmegaConf.load(branch["config_path"])
        model = build_model(cfg).cpu().eval()
        observed = {}

        def record_projection(_module, inputs, output):
            observed["projection_input"] = list(inputs[0].shape)
            observed["projection_output"] = list(output.shape)

        encoder_hook = model.stft_encoder.register_forward_hook(
            lambda _module, _inputs, output, store=observed: store.__setitem__("encoder", list(output.shape))
        )
        projection_hook = model.stft_proj.register_forward_hook(record_projection)
        with torch.no_grad():
            output = model(torch.zeros(1, 1, 18_000))
        encoder_hook.remove()
        projection_hook.remove()

        assert model.fusion_mode == "native_inject"
        assert model.stft_inject_position == "pre_mixer"
        assert model.stft_proj.kernel_size == (1,)
        assert torch.count_nonzero(model.stft_proj.weight) == 0
        assert torch.count_nonzero(model.stft_proj.bias) == 0
        assert observed["encoder"] == [1, 16, branch["raw_stft_shape"][-1]]
        assert observed["projection_input"] == [1, 16, 140]
        assert observed["projection_output"] == [1, 16, 140]
        assert list(output.shape) == [1, 1, 18_000]
        assert branch["full_model_output_shape"] == ["B", 1, 18_000]
        assert branch["fusion_mode"] == "native_inject"
        assert branch["projection_kernel_size"] == [1]
        assert branch["projection_zero_initialized"] is True
