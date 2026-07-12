from __future__ import annotations

from collections.abc import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BODY_BLACK = RGBColor(0x00, 0x00, 0x00)
SCNU_BLUE = RGBColor(0x00, 0x3F, 0x73)
SECONDARY_ORANGE = RGBColor(0xD9, 0x78, 0x24)
NOTE_GRAY = RGBColor(0x66, 0x66, 0x66)
PLACEHOLDER_YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_run_font(run, *, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_shape_text(
    shape,
    text: str,
    *,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    _set_run_font(run, size=size, color=color, bold=bold)


def new_content_slide(prs: Presentation, title: str, *, page_number: int):
    content_layout = next(
        (layout for layout in prs.slide_layouts if layout.name in {"两栏内容", "Two Content"}),
        prs.slide_layouts[4],
    )
    slide = prs.slides.add_slide(content_layout)
    for placeholder in list(slide.placeholders):
        element = placeholder._element
        element.getparent().remove(element)
    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.42), Inches(11.85), Inches(0.58))
    title_box.name = "页面标题"
    _set_shape_text(title_box, title, size=28, color=SCNU_BLUE, bold=True)

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.08), Inches(11.85), Pt(1.2))
    rule.name = "页眉分隔线"
    rule.fill.solid()
    rule.fill.fore_color.rgb = SCNU_BLUE
    rule.line.fill.background()

    page_box = slide.shapes.add_textbox(Inches(12.20), Inches(7.12), Inches(0.45), Inches(0.20))
    page_box.name = "页码"
    _set_shape_text(page_box, str(page_number), size=10, color=NOTE_GRAY, align=PP_ALIGN.RIGHT)
    return slide


def add_text_box(
    slide,
    text: str,
    *,
    x,
    y,
    width,
    height,
    font_size: int = 20,
    color: RGBColor = BODY_BLACK,
    bold: bool = False,
    name: str = "正文",
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(x, y, width, height)
    shape.name = name
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_shape_text(shape, text, size=font_size, color=color, bold=bold, align=align)
    return shape


def add_takeaway(slide, text: str, *, y=Inches(1.23)):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.76), y, Inches(11.78), Inches(0.58))
    bar.name = "核心结论"
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xEC, 0xF3, 0xF8)
    bar.line.color.rgb = SCNU_BLUE
    bar.line.width = Pt(1)
    _set_shape_text(bar, text, size=18, color=BODY_BLACK, bold=True, align=PP_ALIGN.CENTER)
    return bar


def add_source_note(slide, sources: Sequence[str]) -> None:
    text = "来源：" + "；".join(sources)
    add_text_box(
        slide,
        text,
        x=Inches(0.76),
        y=Inches(7.02),
        width=Inches(11.2),
        height=Inches(0.24),
        font_size=9,
        color=NOTE_GRAY,
        name="来源注释",
    )


def add_placeholder(slide, text: str, *, x, y, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.name = "可编辑占位符"
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_YELLOW
    shape.fill.transparency = 25
    shape.line.color.rgb = RGBColor(0xC9, 0x9A, 0x00)
    shape.line.dash_style = 2
    _set_shape_text(shape, text, size=16, color=BODY_BLACK, bold=False, align=PP_ALIGN.CENTER)
    return shape


def add_card(
    slide,
    title: str,
    body: str,
    *,
    x,
    y,
    width,
    height,
    accent: RGBColor = SCNU_BLUE,
    name: str = "信息卡片",
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.4)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(10)
    frame.margin_right = Pt(10)
    frame.margin_top = Pt(8)
    frame.margin_bottom = Pt(8)
    title_paragraph = frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.add_run()
    title_run.text = title
    _set_run_font(title_run, size=18, color=accent, bold=True)
    body_paragraph = frame.add_paragraph()
    body_paragraph.space_before = Pt(8)
    body_paragraph.alignment = PP_ALIGN.LEFT
    body_run = body_paragraph.add_run()
    body_run.text = body
    _set_run_font(body_run, size=14, color=BODY_BLACK)
    return shape


def add_body_text(
    slide,
    paragraphs: Sequence[str],
    *,
    x,
    y,
    width,
    height,
    font_size: int = 20,
    bullet: bool = False,
):
    shape = slide.shapes.add_textbox(x, y, width, height)
    shape.name = "正文"
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, text in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        if bullet:
            paragraph.text = f"• {text}"
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
            run.text = text
        _set_run_font(run, size=font_size, color=BODY_BLACK)
    return shape


def _remove_cell_borders(cell) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr"):
        for node in list(tc_pr.findall(qn(f"a:{edge}"))):
            tc_pr.remove(node)


def _add_cell_border(cell, edge: str, *, width: int = 12700) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    line = OxmlElement(f"a:{edge}")
    line.set("w", str(width))
    solid = OxmlElement("a:solidFill")
    color = OxmlElement("a:srgbClr")
    color.set("val", "000000")
    solid.append(color)
    line.append(solid)
    dash = OxmlElement("a:prstDash")
    dash.set("val", "solid")
    line.append(dash)
    tc_pr.append(line)


def add_three_line_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[object]],
    *,
    x,
    y,
    width,
    height,
    bold_cells: Iterable[tuple[int, int]] = (),
    font_size: int = 16,
):
    table = slide.shapes.add_table(1 + len(rows), len(headers), x, y, width, height).table
    table.first_row = False
    table.first_col = False
    table.last_row = False
    table.last_col = False
    table.horz_banding = False
    table.vert_banding = False
    bold_cells = set(bold_cells)
    values = [list(headers), *[list(row) for row in rows]]
    for row_index, row_values in enumerate(values):
        for column_index, value in enumerate(row_values):
            cell = table.cell(row_index, column_index)
            _remove_cell_borders(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.text = str(value)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            _set_run_font(
                run,
                size=font_size,
                color=BODY_BLACK,
                bold=row_index == 0 or (row_index, column_index) in bold_cells,
            )

    last_row = len(rows)
    for column_index in range(len(headers)):
        _add_cell_border(table.cell(0, column_index), "lnT", width=19050)
        _add_cell_border(table.cell(0, column_index), "lnB", width=12700)
        _add_cell_border(table.cell(last_row, column_index), "lnB", width=19050)
    return table


def is_three_line_table(table) -> bool:
    rows = len(table.rows)
    columns = len(table.columns)
    if rows < 2 or columns < 1:
        return False
    for row_index in range(rows):
        for column_index in range(columns):
            tc_pr = table.cell(row_index, column_index)._tc.get_or_add_tcPr()
            present = {
                edge
                for edge in ("lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr")
                if tc_pr.find(qn(f"a:{edge}")) is not None
            }
            expected = set()
            if row_index == 0:
                expected.update({"lnT", "lnB"})
            if row_index == rows - 1:
                expected.add("lnB")
            if present != expected:
                return False
    return True
