"""研究讨论型正文的显式页组与原生 PowerPoint builders。"""

from __future__ import annotations

import csv
import hashlib
import json
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Mapping

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from .discussion_content import DISCUSSION_UNITS, DiscussionUnit
from .evidence import FORMAL_LABELS, build_evidence_catalog
from .theme import (
    SCNU_BLUE,
    SECONDARY_ORANGE,
    WHITE,
    add_discussion_box,
    add_evidence_boundary,
    add_evidence_gap,
    add_method_panel,
    add_source_note,
    add_text_box,
    add_three_line_table,
    new_content_slide,
)


@dataclass(frozen=True)
class PanelPlan:
    field: str
    title: str
    role: str = "method"


@dataclass(frozen=True)
class SlideSpec:
    """人工审查后的页定义；内容选择和版面均由 spec 显式声明。"""

    key: str
    section: str
    title: str
    unit_keys: tuple[str, ...]
    context_keys: tuple[str, ...]
    builder: str
    asset_key: str | None
    display_fields: tuple[str, ...]
    panel_plan: tuple[PanelPlan, ...]


QUESTION = "question"
METHOD = "method_steps"
PARAMETERS = "parameters"
RATIONALE = "rationale"
EVIDENCE = "evidence"
LIMITS = "limits"
PROMPT = "discussion_prompt"


def _p(field: str, title: str, role: str = "method") -> PanelPlan:
    return PanelPlan(field, title, role)


def _s(
    key: str,
    section: str,
    title: str,
    *unit_keys: str,
    builder: str,
    fields: tuple[str, ...],
    plan: tuple[PanelPlan, ...],
    asset_key: str | None = None,
    context_keys: tuple[str, ...] = (),
) -> SlideSpec:
    if QUESTION not in fields:
        raise ValueError(f"SlideSpec {key} 必须显式展示 question")
    planned = {panel.field for panel in plan}
    if planned != set(fields) - {QUESTION}:
        raise ValueError(f"SlideSpec {key} 的 display_fields 与 panel_plan 不一致")
    return SlideSpec(key, section, title, tuple(unit_keys), context_keys, builder, asset_key, fields, plan)


GENERIC_METHOD = (
    (QUESTION, METHOD, PARAMETERS, PROMPT),
    (_p(METHOD, "具体步骤"), _p(PARAMETERS, "参数 / shape"), _p(PROMPT, "需要建议 / 待决策", "discussion")),
)
GENERIC_REASON = (
    (QUESTION, RATIONALE, LIMITS, PROMPT),
    (_p(RATIONALE, "为什么"), _p(LIMITS, "限制", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")),
)
ASSET_PLAN = (
    (QUESTION, PARAMETERS, LIMITS, PROMPT),
    (_p(PARAMETERS, "步骤与参数"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")),
)


# continuation 页通过 context_keys 复用语义，但 48 个 unit key 仍恰好消费一次。
SLIDE_GROUPS = (
    _s("task_observation", "任务与信号直觉", "从混合观测到 THO-like 呼吸表示", "task_mapping", "bcg_signal_mixture", builder="asset", fields=(QUESTION, RATIONALE, LIMITS, PROMPT), plan=(_p(RATIONALE, "为什么"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="signal_overview"),
    _s("reference_boundary", "任务与信号直觉", "参考信号也有边界：时延、方向与局部质量", "tho_reference_limits", builder="generic", fields=GENERIC_REASON[0], plan=GENERIC_REASON[1]),
    _s("window_timescale", "任务与信号直觉", "三分钟样本如何兼顾稳定估计与局部变化", "window_180s_tradeoff", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),

    _s("data_provenance", "数据来源与样本形成", "数据口径为何必须由导出 provenance 冻结", "data_provenance", builder="provenance", fields=(QUESTION, METHOD, RATIONALE, LIMITS, PROMPT), plan=(_p(METHOD, "追溯步骤"), _p(RATIONALE, "为什么"), _p(LIMITS, "当前边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("sample_formation", "数据来源与样本形成", "样本如何从整晚数据形成", "index_to_window", builder="sample", fields=(QUESTION, METHOD, PARAMETERS, RATIONALE, LIMITS), plan=(_p(METHOD, "四步形成样本"), _p(PARAMETERS, "参数 / shape"), _p(RATIONALE, "为什么"), _p(LIMITS, "限制", "evidence"))),
    _s("window_filter", "数据来源与样本形成", "进入模型前：waveform 窗口筛选", "waveform_window_filter", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),
    _s("split_boundary", "数据来源与样本形成", "受试者隔离定义数据泄漏边界", "subject_split_boundary", builder="generic", fields=GENERIC_REASON[0], plan=GENERIC_REASON[1]),

    _s("wideband_input", "输入与目标预处理", "为什么保留 0.03–20 Hz 宽频 BCG", "bcg_wideband_input", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),
    _s("robust_z", "输入与目标预处理", "robust-z：按状态段统一尺度", "bcg_segment_robust_z", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),
    _s("softz", "输入与目标预处理", "soft-z：压缩长尾，但本窗未触发压缩", "bcg_soft_z", builder="softz", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="softz_mapping"),
    _s("target_pipeline", "输入与目标预处理", "目标、整窗筛选与 RR mask：三条链路", "target_soft_z_and_mask", builder="target_pipeline", fields=(QUESTION, METHOD, PARAMETERS), plan=(_p(METHOD, "三条链路"), _p(PARAMETERS, "冻结参数"))),
    _s("target_evidence", "输入与目标预处理", "目标监督代码证据：完整 target 与评价 mask", builder="evidence_page", fields=(QUESTION, EVIDENCE), plan=(_p(EVIDENCE, "代码证据", "evidence"),), context_keys=("target_soft_z_and_mask",)),
    _s("target_boundary", "输入与目标预处理", "目标监督边界：训练损失与 RR mask 不混用", builder="generic", fields=(QUESTION, RATIONALE, LIMITS, PROMPT), plan=(_p(RATIONALE, "为什么必须区分"), _p(LIMITS, "当前限制", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), context_keys=("target_soft_z_and_mask",)),

    _s("tensor_contract", "模型输入与计算图", "模型张量从 B×1×18000 开始", "model_tensor_contract", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),
    _s("patch_geometry", "模型输入与计算图", "PatchMixer：256 / 128 / 140 的张量契约", "patchmixer_token_shape", builder="asset", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="token_geometry"),
    _s("stft_shapes", "模型输入与计算图", "STFT：F0 / wide / bandenergy shape", "stft_resolution", builder="asset", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="stft_branch_shapes"),
    _s("fusion", "模型输入与计算图", "pre-mixer 加法融合如何进入计算图", "pre_mixer_injection", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),
    _s("bandenergy", "模型输入与计算图", "五条重叠频带的表示能力与 tradeoff", "bandenergy_encoder", builder="asset", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="bandenergy_response"),

    _s("dataset_roles", "训练与损失", "训练、验证与独立测试各自决定什么", "train_val_test_roles", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),
    _s("loss_goal", "训练与损失", "WeakSyncLoss：为什么不是单一逐点误差", "composite_loss_goal", builder="loss_goal", fields=(QUESTION, METHOD, RATIONALE, LIMITS, PROMPT), plan=(_p(METHOD, "组合步骤"), _p(RATIONALE, "为什么"), _p(LIMITS, "限制", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("loss_weights", "训练与损失", "WeakSyncLoss：五项基础权重与任务含义", "loss_component_weights", builder="loss_weights", fields=(QUESTION, METHOD, PARAMETERS, RATIONALE, LIMITS, PROMPT), plan=(_p(METHOD, "损失项"), _p(PARAMETERS, "权重"), _p(RATIONALE, "为什么"), _p(LIMITS, "限制", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("signed_loss", "训练与损失", "L_signed_corr 与方向调度", "polarity_schedule", builder="asset_formula", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="loss_schedule"),
    _s("checkpoint_flow", "训练与损失", "checkpoint：方向 gate 与 val-loss top-k", "checkpoint_selection", builder="checkpoint_flow", fields=(QUESTION, METHOD, PARAMETERS), plan=(_p(METHOD, "选择流程"), _p(PARAMETERS, "冻结参数"))),
    _s("checkpoint_evidence", "训练与损失", "checkpoint：canonical test 旁路复评证据链", builder="evidence_page", fields=(QUESTION, EVIDENCE), plan=(_p(EVIDENCE, "证据链", "evidence"),), context_keys=("checkpoint_selection",)),
    _s("checkpoint_boundary", "训练与损失", "checkpoint：task checkpoint 与 val-loss 语义边界", builder="generic", fields=(QUESTION, LIMITS, PROMPT), plan=(_p(LIMITS, "语义边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), context_keys=("checkpoint_selection",)),

    _s("robust_rr", "指标计算与失效场景", "稳健 RR：谱引导、峰间距与双峰护栏", "metric_robust_rr", builder="asset_formula", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="metric_robust_rr"),
    _s("cycle", "指标计算与失效场景", "cycle：双向过零如何变成周期计数", "metric_breath_count", builder="asset_formula", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="metric_cycle_count"),
    _s("lag4", "指标计算与失效场景", "lag4：在 ±4 秒内分离形态与时延", "metric_lag_corr", builder="asset_formula", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="metric_lag_corr"),
    _s("relative_envelope", "指标计算与失效场景", "relative envelope：比较局部增强与减弱", "metric_relative_envelope", builder="asset_formula", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="metric_relative_envelope"),
    _s("local_rr", "指标计算与失效场景", "local RR：正式 40 秒窗 / 10 秒步长", "metric_local_rr", builder="local_rr", fields=(QUESTION, METHOD, PARAMETERS, LIMITS, PROMPT), plan=(_p(METHOD, "完整步骤"), _p(PARAMETERS, "正式参数"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="metric_local_rr"),

    _s("controls", "对照实验设计", "四个对照只改变时频辅助表示", "controlled_variables", builder="controls", fields=(QUESTION, LIMITS), plan=(_p(LIMITS, "对照边界", "evidence"),)),
    _s("time_only", "对照实验设计", "time-only substrate 的计算边界", "time_only_baseline", builder="generic", fields=GENERIC_REASON[0], plan=GENERIC_REASON[1]),
    _s("f0_wide", "对照实验设计", "F0 与 wide：分辨力和更新速度的取舍", "f0_vs_wide", builder="asset", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="stft_resolution_comparison"),
    _s("wide_band", "对照实验设计", "完整谱图还是低维频带摘要", "wide_vs_bandenergy", builder="stft_table", fields=(QUESTION, RATIONALE, LIMITS, PROMPT), plan=(_p(RATIONALE, "为什么"), _p(LIMITS, "限制", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),

    _s("overall_primary", "整体结果与稳定性", "四方案统一比较（一）：整窗节律与形态", "overall_metric_values", builder="canonical_primary", fields=(QUESTION, EVIDENCE, LIMITS), plan=(_p(EVIDENCE, "canonical evidence", "evidence"), _p(LIMITS, "统计边界", "evidence"))),
    _s("overall_secondary", "整体结果与稳定性", "四方案统一比较（二）：相对包络与 local RR", builder="canonical_secondary", fields=(QUESTION, PROMPT), plan=(_p(PROMPT, "需要建议 / 待决策", "discussion"),), context_keys=("overall_metric_values",)),
    _s("paired", "整体结果与稳定性", "paired delta：同 seed 隔离方案差异", "paired_delta", builder="asset", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="overall_delta"),
    _s("seed_subject", "整体结果与稳定性", "seed × subject：方向稳定不等于总体稳定", "seed_stability", builder="asset", fields=ASSET_PLAN[0], plan=ASSET_PLAN[1], asset_key="seed_subject_stability"),
    _s("decision_rule", "整体结果与稳定性", "结论顺序：主护栏、形态、局部诊断", "result_decision_rule", builder="generic", fields=GENERIC_REASON[0], plan=GENERIC_REASON[1]),

    _s("subject", "分层分析与完整案例", "受试者分层：总体均值由谁贡献", "subject_stratification", builder="generic", fields=(QUESTION, METHOD, RATIONALE, PROMPT), plan=(_p(METHOD, "具体步骤"), _p(RATIONALE, "为什么"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("quality", "分层分析与完整案例", "质量与难度分层：收益发生在哪里", "quality_difficulty_strata", builder="asset", fields=(QUESTION, RATIONALE, LIMITS, PROMPT), plan=(_p(RATIONALE, "为什么"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="strata_tradeoffs"),
    _s("rr_bins", "分层分析与完整案例", "RR 分层：慢、常规与快呼吸的角色差异", "rr_bin_strata", builder="generic", fields=(QUESTION, METHOD, RATIONALE, PROMPT), plan=(_p(METHOD, "具体步骤"), _p(RATIONALE, "为什么"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("case_640", "分层分析与完整案例", "dataset_row_id=640｜四模型均纠正", "complete_case_review", builder="case", fields=(QUESTION, METHOD, LIMITS, PROMPT), plan=(_p(METHOD, "完整复核"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="case_row_640"),
    _s("case_873", "分层分析与完整案例", "dataset_row_id=873｜四模型均失败", builder="case", fields=(QUESTION, METHOD, LIMITS, PROMPT), plan=(_p(METHOD, "完整复核"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="case_row_873", context_keys=("complete_case_review",)),
    _s("case_1353", "分层分析与完整案例", "dataset_row_id=1353｜模型间分歧", builder="case", fields=(QUESTION, METHOD, LIMITS, PROMPT), plan=(_p(METHOD, "完整复核"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="case_row_1353", context_keys=("complete_case_review",)),
    _s("case_3584", "分层分析与完整案例", "dataset_row_id=3584｜阈值边界", builder="case", fields=(QUESTION, METHOD, LIMITS, PROMPT), plan=(_p(METHOD, "完整复核"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion")), asset_key="case_row_3584", context_keys=("complete_case_review",)),

    _s("harmonic_definition", "BCG 二次谐波问题", "二次谐波：BCG 主峰为什么落在 2×THO", "harmonic_definition", builder="formula", fields=(QUESTION, METHOD, PARAMETERS, LIMITS, PROMPT), plan=(_p(METHOD, "判定步骤"), _p(PARAMETERS, "参数"), _p(LIMITS, "限制", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("harmonic_threshold", "BCG 二次谐波问题", "二次谐波判据：阈值如何冻结", "harmonic_thresholds", builder="formula", fields=(QUESTION, METHOD, PARAMETERS, LIMITS, PROMPT), plan=(_p(METHOD, "冻结步骤"), _p(PARAMETERS, "阈值"), _p(LIMITS, "边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("harmonic_types_boundary", "BCG 二次谐波问题", "三类谐波标签与边界案例如何共同解释", "harmonic_subtypes", "harmonic_boundary_cases", builder="generic", fields=(QUESTION, RATIONALE, LIMITS, PROMPT), plan=(_p(RATIONALE, "为什么分开"), _p(LIMITS, "外推边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("harmonic_correction", "BCG 二次谐波问题", "纠偏：回到基频且压低相对二倍频", "harmonic_correction", builder="generic", fields=GENERIC_METHOD[0], plan=GENERIC_METHOD[1]),

    _s("output_and_gate", "研究议题与下一步", "输出空间与 BCG-only gate 必须一起决策", "decision_output_space", "decision_bcg_only_gate", builder="decision_pair", fields=(QUESTION, RATIONALE, LIMITS, PROMPT), plan=(_p(RATIONALE, "为什么"), _p(LIMITS, "证据边界", "evidence"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("stat_unit", "研究议题与下一步", "统计单位从重叠窗提升到受试者与时间块", "decision_statistical_unit", builder="generic", fields=(QUESTION, METHOD, RATIONALE, PROMPT), plan=(_p(METHOD, "具体步骤"), _p(RATIONALE, "为什么"), _p(PROMPT, "需要建议 / 待决策", "discussion"))),
    _s("minimal_next", "研究议题与下一步", "下一轮最小判别性实验", "decision_minimal_next_experiment", builder="minimal_next", fields=(QUESTION, METHOD, RATIONALE, LIMITS, PROMPT), plan=(_p(METHOD, "最小实验步骤"), _p(RATIONALE, "为什么"), _p(LIMITS, "停止边界", "evidence"), _p(PROMPT, "待决策", "discussion"))),
)


UNIT_BY_KEY = {unit.key: unit for unit in DISCUSSION_UNITS}
MODEL_NAMES = {
    "g0_time_only": "纯时序",
    "g0_f0_native_stft_pre_mixer": "F0 STFT",
    "g3_c_wide_8p0": "wide STFT",
    "g3_c_bandenergy": "bandenergy",
}
ASSET_FILES = {
    key: f"{key}.png" for key in (
        "signal_overview", "softz_mapping", "token_geometry", "stft_branch_shapes",
        "bandenergy_response", "loss_schedule", "metric_robust_rr", "metric_cycle_count",
        "metric_lag_corr", "metric_relative_envelope", "metric_local_rr", "overall_delta",
        "seed_subject_stability", "strata_tradeoffs", "stft_resolution_comparison",
        "case_row_640", "case_row_873", "case_row_1353", "case_row_3584",
    )
}


_MANIFEST_ASSETS = {
    "signal_assets_manifest.json": {
        "signal_overview", "softz_mapping", "stft_resolution_comparison", "bandenergy_response",
    },
    "model_assets_manifest.json": {"token_geometry", "stft_branch_shapes", "loss_schedule"},
    "metric_assets_manifest.json": {
        "metric_robust_rr", "metric_cycle_count", "metric_lag_corr",
        "metric_relative_envelope", "metric_local_rr",
    },
    "case_assets_manifest.json": {
        "overall_delta", "seed_subject_stability", "strata_tradeoffs",
        "case_row_640", "case_row_873", "case_row_1353", "case_row_3584",
    },
}


def _record_path(repo_root: Path, record: Mapping[str, object]) -> Path:
    path = Path(str(record.get("path", "")))
    return path if path.is_absolute() else repo_root / path


def _validate_record(manifest_path: Path, label: str, record: object, actual_path: Path) -> None:
    if not isinstance(record, Mapping):
        raise ValueError(f"记录 {label} 缺失或格式错误")
    if not actual_path.is_file():
        raise ValueError(f"记录 {label} 指向的文件不存在: {actual_path}")
    size = actual_path.stat().st_size
    expected_size = record.get("size_bytes")
    expected_sha = record.get("sha256")
    if expected_size != size:
        raise ValueError(f"记录 {label} 大小不匹配: manifest={expected_size}, actual={size}")
    digest = hashlib.sha256()
    with actual_path.open("rb") as fp:
        for chunk in iter(lambda: fp.read(1024 * 1024), b""):
            digest.update(chunk)
    if expected_sha != digest.hexdigest():
        raise ValueError(f"记录 {label} SHA256 不匹配: {actual_path}")


def _validate_manifest_sources(repo_root: Path, manifest_name: str, manifest: Mapping[str, object]) -> None:
    if manifest_name == "signal_assets_manifest.json":
        evidence = manifest.get("evidence", {})
        configs = manifest.get("resolved_configs", {})
        # checkpoint、原始数据文件不参与组装期哈希；这里只校验可快速复核的小型关键来源。
        for key in ("signal_figure_code", "signal_metadata_json", "signal_npz", "source_run_config"):
            record = evidence.get(key) if isinstance(evidence, Mapping) else None
            _validate_record(Path(manifest_name), f"evidence.{key}", record, _record_path(repo_root, record or {}))
        for key in ("f0", "wide", "bandenergy"):
            record = configs.get(key) if isinstance(configs, Mapping) else None
            _validate_record(Path(manifest_name), f"resolved_configs.{key}", record, _record_path(repo_root, record or {}))
    elif manifest_name == "model_assets_manifest.json":
        sources = manifest.get("sources", {})
        for key in ("model_figure_code", "patch_mixer_code", "stft_branch_code", "loss_code"):
            record = sources.get(key) if isinstance(sources, Mapping) else None
            _validate_record(Path(manifest_name), f"sources.{key}", record, _record_path(repo_root, record or {}))
        configs = sources.get("formal_configs", {}) if isinstance(sources, Mapping) else {}
        for key, record in configs.items() if isinstance(configs, Mapping) else ():
            _validate_record(Path(manifest_name), f"sources.formal_configs.{key}", record, _record_path(repo_root, record))
    elif manifest_name == "metric_assets_manifest.json":
        sources = manifest.get("sources", {})
        for key in (
            "signal_metadata_json", "source_run_config", "source_metrics_csv", "signal_npz",
            "metric_code", "metric_evaluate_code", "metric_figure_code", "robust_rr_demo_code", "metric_schema",
        ):
            record = sources.get(key) if isinstance(sources, Mapping) else None
            _validate_record(Path(manifest_name), f"sources.{key}", record, _record_path(repo_root, record or {}))
    else:
        if manifest.get("generator") != "scripts.tho_group_meeting_ppt.case_figures.build_case_assets":
            raise ValueError("generator 标识错误")
        sources = manifest.get("sources", ())
        selected = []
        for record in sources if isinstance(sources, list) else ():
            path = Path(str(record.get("path", "")))
            if path.is_absolute():
                continue
            name = path.name
            if name == "case_figures.py" or name.endswith(("_manifest.json", "_manifest.csv", "_summary.csv", "_labels.csv")) or name == "harmonic_thresholds.json":
                selected.append(record)
        if not any(Path(str(record.get("path"))).name == "case_figures.py" for record in selected):
            raise ValueError("缺少 case_figures.py 生成器源码记录")
        for record in selected:
            _validate_record(Path(manifest_name), f"sources.{record.get('path')}", record, _record_path(repo_root, record))


def resolve_discussion_assets(repo_root: Path, asset_dir: Path | None = None) -> dict[str, Path]:
    """读取四份 manifest，并在组装前验证 19 张图及关键生成证据。"""
    root = asset_dir or repo_root / "docs/stage_reports/20260708/generated_assets/discussion"
    validated: dict[str, Path] = {}
    for manifest_name, expected_keys in _MANIFEST_ASSETS.items():
        manifest_path = root / manifest_name
        try:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            records = manifest.get("assets", {})
            missing = expected_keys - set(records) if isinstance(records, Mapping) else expected_keys
            if missing:
                raise ValueError(f"缺少资产记录: {sorted(missing)}")
            for key in expected_keys:
                path = root / ASSET_FILES[key]
                _validate_record(manifest_path, f"assets.{key}", records[key], path)
                validated[key] = path
            _validate_manifest_sources(repo_root, manifest_name, manifest)
        except (OSError, json.JSONDecodeError, ValueError, TypeError) as exc:
            raise ValueError(
                f"{manifest_name} 证据校验失败: {exc}；请先运行 --assets-only 重新生成讨论资产"
            ) from exc
    if set(validated) != set(ASSET_FILES):
        raise ValueError("讨论资产集合不是预期的 19 项；请先运行 --assets-only 重新生成讨论资产")
    return validated


def load_canonical_overall_metrics(repo_root: Path) -> tuple[dict[str, float | str], ...]:
    """从 EvidenceCatalog 已验证 manifest 所属目录读取 canonical 四方案汇总。"""
    catalog = build_evidence_catalog(repo_root)
    path = catalog.result_root / "g_series_local_rr_canonical_label_summary.csv"
    required = {
        "label", "rr_peak_band_robust_abs_error_mean_mean",
        "breath_count_zero_cross_abs_error_mean_mean", "best_lag_corr_4s_mean_mean",
        "relative_envelope_corr_lag4s_mean_mean", "local_rr_mae_mean_mean",
    }
    with path.open(newline="", encoding="utf-8") as fp:
        reader = csv.DictReader(fp)
        missing = required - set(reader.fieldnames or ())
        if missing:
            raise ValueError(f"canonical summary 缺少列 {sorted(missing)}: {path}")
        by_label = {row["label"]: row for row in reader}
    if set(FORMAL_LABELS) - set(by_label):
        raise ValueError(f"canonical summary 缺少正式方案: {path}")
    rows = []
    for label in FORMAL_LABELS:
        row = by_label[label]
        rows.append({
            "label": label,
            "robust_rr_mae_bpm": float(row["rr_peak_band_robust_abs_error_mean_mean"]),
            "count_bpm_mae": float(row["breath_count_zero_cross_abs_error_mean_mean"]) / 3.0,
            "lag4_corr": float(row["best_lag_corr_4s_mean_mean"]),
            "relative_envelope_corr": float(row["relative_envelope_corr_lag4s_mean_mean"]),
            "local_rr_mae_bpm": float(row["local_rr_mae_mean_mean"]),
        })
    return tuple(rows)


def _units(spec: SlideSpec) -> tuple[DiscussionUnit, ...]:
    keys = spec.unit_keys or spec.context_keys
    return tuple(UNIT_BY_KEY[key] for key in keys)


def _field_text(units: tuple[DiscussionUnit, ...], field: str) -> str:
    if field == QUESTION:
        values = [unit.question for unit in units]
    elif field == PROMPT:
        values = [unit.discussion_prompt for unit in units if unit.discussion_prompt]
    else:
        values = [value for unit in units for value in getattr(unit, field)]
    cleaned = [str(value).replace("更优", "更合适") for value in values]
    if field in {QUESTION, PROMPT}:
        return " / ".join(cleaned)
    return "\n".join(f"• {value}" for value in cleaned)


def _picture(slide, path: Path, *, x=Inches(0.65), y=Inches(2.00), width=Inches(6.20), height=Inches(3.50)):
    with Image.open(path) as image:
        ratio = image.width / image.height
        source: str | BytesIO = str(path)
        if path.name.startswith("case_row_") and image.mode == "RGBA":
            # LibreOffice 对部分全不透明 RGBA 案例图会生成黑色透明组；嵌入前仅移除
            # 无信息量 alpha 通道，像素 RGB 与源图保持一致。
            source = BytesIO()
            image.convert("RGB").save(source, format="PNG")
            source.seek(0)
    target_ratio = width / height
    if ratio > target_ratio:
        draw_width, draw_height = width, int(width / ratio)
    else:
        draw_height, draw_width = height, int(height * ratio)
    pic = slide.shapes.add_picture(source, int(x + (width - draw_width) / 2), int(y + (height - draw_height) / 2), draw_width, draw_height)
    pic.name = f"真实证据图：{path.name}"
    return pic


def _base(prs: Presentation, spec: SlideSpec, page: int, units: tuple[DiscussionUnit, ...]):
    slide = new_content_slide(prs, spec.title, page_number=page)
    add_text_box(slide, f"研究问题｜{_field_text(units, QUESTION)}", x=Inches(0.76), y=Inches(1.23), width=Inches(11.10), height=Inches(0.56), font_size=18, bold=True, name="正文研究问题", align=PP_ALIGN.CENTER)
    add_source_note(slide, tuple(dict.fromkeys(source for unit in units for source in unit.sources)))
    return slide


def _add_planned_panel(slide, plan: PanelPlan, body: str, *, x, y, width, height):
    if plan.role == "discussion":
        return add_discussion_box(slide, body, x=x, y=y, width=width, height=height)
    if plan.role == "evidence":
        return add_evidence_boundary(slide, body, x=x, y=y, width=width, height=height)
    return add_method_panel(slide, plan.title, body, x=x, y=y, width=width, height=height)


def _render_generic(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    discussion = [panel for panel in spec.panel_plan if panel.role == "discussion"]
    content = [panel for panel in spec.panel_plan if panel.role != "discussion"]
    if len(discussion) > 1 or len(content) > 3:
        raise ValueError(f"SlideSpec {spec.key} 的通用布局面板过多")
    bottom = Inches(5.38) if discussion else Inches(6.48)
    height = bottom - Inches(2.05) - Inches(0.20)
    gap = Inches(0.20)
    width = int((Inches(11.82) - gap * (len(content) - 1)) / max(1, len(content)))
    for index, panel in enumerate(content):
        _add_planned_panel(slide, panel, _field_text(units, panel.field), x=Inches(0.72) + index * (width + gap), y=Inches(2.05), width=width, height=height)
    if discussion:
        panel = discussion[0]
        _add_planned_panel(slide, panel, _field_text(units, panel.field), x=Inches(1.05), y=Inches(5.43), width=Inches(11.15), height=Inches(1.18))


def _native_flow(slide, steps: tuple[str, ...], *, y: float = 2.05, height: float = 1.25):
    count = len(steps)
    if not 1 <= count <= 4:
        raise ValueError("专用流程页仅支持显式的 1–4 步")
    gap = 0.32
    width = (11.85 - gap * (count - 1)) / count
    for index, text in enumerate(steps):
        x = 0.72 + index * (width + gap)
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        node.name = "流程节点"
        node.fill.solid(); node.fill.fore_color.rgb = WHITE; node.line.color.rgb = SCNU_BLUE
        add_text_box(slide, text, x=Inches(x + 0.10), y=Inches(y + 0.08), width=Inches(width - 0.20), height=Inches(height - 0.16), font_size=18, name="流程节点文字", align=PP_ALIGN.CENTER)
        if index:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x - 0.28), Inches(y + height / 2), Inches(x - 0.04), Inches(y + height / 2))
            line.name = "流程连接线"; line.line.color.rgb = SECONDARY_ORANGE; line.line.end_arrowhead = True


def _render_asset(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...], path: Path):
    _picture(slide, path)
    discussion = [panel for panel in spec.panel_plan if panel.role == "discussion"]
    content = [panel for panel in spec.panel_plan if panel.role != "discussion"]
    if len(content) != 2 or len(discussion) != 1:
        raise ValueError(f"SlideSpec {spec.key} 的资产布局必须为 2 个内容面板 + 1 个讨论框")
    _add_planned_panel(slide, content[0], _field_text(units, content[0].field), x=Inches(7.10), y=Inches(2.00), width=Inches(5.15), height=Inches(1.72))
    boundary = _field_text(units, content[1].field)
    boundary_y = Inches(3.92)
    boundary_height = Inches(1.35)
    if spec.builder == "softz":
        boundary = "• 压缩改变幅值物理意义，可能弱化真实大呼吸。\n• 本窗 0 samples changed；不外推其他窗口。"
        boundary_y = Inches(3.74)
        boundary_height = Inches(1.68)
    _add_planned_panel(slide, content[1], boundary, x=Inches(7.10), y=boundary_y, width=Inches(5.15), height=boundary_height)
    _add_planned_panel(slide, discussion[0], _field_text(units, discussion[0].field), x=Inches(0.90), y=Inches(5.48), width=Inches(11.45), height=Inches(1.18))


def _render_provenance(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    content = [panel for panel in spec.panel_plan if panel.role != "discussion"]
    gap = Inches(0.20)
    width = int((Inches(11.82) - gap * (len(content) - 1)) / len(content))
    for index, panel in enumerate(content):
        _add_planned_panel(
            slide,
            panel,
            _field_text(units, panel.field),
            x=Inches(0.72) + index * (width + gap),
            y=Inches(2.05),
            width=width,
            height=Inches(3.05),
        )
    add_evidence_gap(
        slide,
        missing_fields="dirty=false 的数据制作 commit 与重导出 manifest",
        suggested_entrypoint=(
            "在数据制作工作区重新导出并冻结 provenance；待决策："
            + _field_text(units, PROMPT)
        ),
        x=Inches(0.95),
        y=Inches(5.28),
        width=Inches(11.35),
        height=Inches(1.52),
    )


def _render_sample(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    _native_flow(slide, tuple(value for unit in units for value in unit.method_steps))
    plans = {panel.field: panel for panel in spec.panel_plan}
    _add_planned_panel(slide, plans[PARAMETERS], _field_text(units, PARAMETERS), x=Inches(0.75), y=Inches(3.55), width=Inches(3.65), height=Inches(1.75))
    _add_planned_panel(slide, plans[RATIONALE], _field_text(units, RATIONALE), x=Inches(4.58), y=Inches(3.55), width=Inches(3.65), height=Inches(1.75))
    _add_planned_panel(slide, plans[LIMITS], _field_text(units, LIMITS), x=Inches(8.41), y=Inches(3.55), width=Inches(3.65), height=Inches(1.75))


def _render_target_pipeline(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    _native_flow(slide, tuple(value for unit in units for value in unit.method_steps), height=2.20)
    add_method_panel(slide, "冻结参数", _field_text(units, PARAMETERS), x=Inches(1.10), y=Inches(4.45), width=Inches(11.10), height=Inches(1.32))
    add_evidence_boundary(slide, "本页只定义三条链路；训练语义与代码证据在下一页展开。", x=Inches(2.10), y=Inches(5.85), width=Inches(9.10), height=Inches(1.00))


FORMULAS = {
    "signed_loss": "L_signed_corr: 0.6 → 0.2；L_signed_cos: 0.1 → 0（epoch 1–6）",
    "robust_rr": "x → bandpass → Welch 主峰 → find_peaks → median(Δt_peak) → 60 / Δt",
    "cycle": "cycle_bpm_error = |N_cycle(ŷ) − N_cycle(y)| / 3 min",
    "lag4": "corr_lag4 = max corr(ŷ(t+τ), y(t)), |τ| ≤ 4 s",
    "relative_envelope": "R(x) = RMS₂ₛ(x) / smooth₂₀ₛ(RMS₂ₛ(x))",
    "harmonic_definition": "f_BCG,peak ≈ 2 × f_THO,fundamental",
    "harmonic_threshold": "ratio₂f/f + peak-distance + valid-spectrum → frozen label",
}


def _render_loss_goal(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    add_text_box(slide, "L = L_env + 0.2L_spec + 0.1L_smooth + 0.2L_high + 0.03L_rel_env + L_signed", x=Inches(0.82), y=Inches(1.98), width=Inches(11.55), height=Inches(0.62), font_size=22, bold=True, name="公式", align=PP_ALIGN.CENTER)
    plans = {panel.field: panel for panel in spec.panel_plan}
    _add_planned_panel(slide, plans[METHOD], _field_text(units, METHOD), x=Inches(0.75), y=Inches(2.88), width=Inches(5.70), height=Inches(2.12))
    _add_planned_panel(slide, plans[RATIONALE], _field_text(units, RATIONALE), x=Inches(6.75), y=Inches(2.88), width=Inches(2.75), height=Inches(2.12))
    _add_planned_panel(slide, plans[LIMITS], _field_text(units, LIMITS), x=Inches(9.75), y=Inches(2.88), width=Inches(2.75), height=Inches(2.12))
    _add_planned_panel(slide, plans[PROMPT], _field_text(units, PROMPT), x=Inches(1.05), y=Inches(5.38), width=Inches(11.15), height=Inches(1.22))


def _render_loss_weights(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    table = add_three_line_table(slide, ("损失项", "权重", "任务含义"), (
        ("L_env", "1.0", "2 秒 RMS 包络相关"), ("L_spec", "0.2", "0.05–0.7 Hz 归一化功率"),
        ("L_smooth", "0.1", "相邻点差分"), ("L_high", "0.2", "0.7 Hz 以上相对能量"),
        ("L_rel_env", "0.03", "局部相对包络"),
    ), x=Inches(0.65), y=Inches(2.02), width=Inches(7.20), height=Inches(3.65), font_size=16)
    table._graphic_frame.name = "WeakSyncLoss 权重表"
    plans = {panel.field: panel for panel in spec.panel_plan}
    _add_planned_panel(slide, plans[RATIONALE], _field_text(units, RATIONALE), x=Inches(8.10), y=Inches(2.02), width=Inches(4.20), height=Inches(1.55))
    _add_planned_panel(slide, plans[LIMITS], _field_text(units, LIMITS), x=Inches(8.10), y=Inches(3.82), width=Inches(4.20), height=Inches(1.55))
    _add_planned_panel(slide, plans[PROMPT], _field_text(units, PROMPT), x=Inches(1.05), y=Inches(5.72), width=Inches(11.15), height=Inches(1.05))


def _render_checkpoint_flow(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    _native_flow(slide, tuple(value for unit in units for value in unit.method_steps), y=1.95, height=2.18)
    add_method_panel(slide, "冻结参数", _field_text(units, PARAMETERS), x=Inches(1.10), y=Inches(4.42), width=Inches(11.10), height=Inches(1.62))


def _render_evidence_page(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    panel = spec.panel_plan[0]
    _add_planned_panel(slide, panel, _field_text(units, panel.field), x=Inches(0.95), y=Inches(2.05), width=Inches(11.35), height=Inches(4.35))


def _render_local_rr(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...], path: Path):
    _picture(slide, path, width=Inches(6.10), height=Inches(3.43))
    add_method_panel(slide, "完整步骤", _field_text(units, METHOD), x=Inches(7.00), y=Inches(1.98), width=Inches(5.30), height=Inches(2.32))
    add_evidence_boundary(slide, "正式：40 秒窗 / 10 秒步长；20 秒局部 RR 窗属于旧/其他流程，勿混用。\n• 非逐呼吸事件匹配；谱峰切换、谐波与低 valid fraction 会失真。", x=Inches(7.00), y=Inches(4.31), width=Inches(5.30), height=Inches(2.02))
    add_text_box(slide, "待决策｜" + _field_text(units, PROMPT), x=Inches(0.90), y=Inches(6.40), width=Inches(11.45), height=Inches(0.42), font_size=18, bold=True, name="讨论框正文", align=PP_ALIGN.CENTER)


def _render_formula(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    add_text_box(slide, FORMULAS[spec.key], x=Inches(0.82), y=Inches(1.98), width=Inches(11.55), height=Inches(0.62), font_size=22, bold=True, name="公式", align=PP_ALIGN.CENTER)
    plans = {panel.field: panel for panel in spec.panel_plan}
    _add_planned_panel(slide, plans[METHOD], _field_text(units, METHOD), x=Inches(0.72), y=Inches(2.88), width=Inches(5.30), height=Inches(2.22))
    _add_planned_panel(slide, plans[PARAMETERS], _field_text(units, PARAMETERS), x=Inches(6.22), y=Inches(2.88), width=Inches(2.95), height=Inches(2.22))
    _add_planned_panel(slide, plans[LIMITS], _field_text(units, LIMITS), x=Inches(9.37), y=Inches(2.88), width=Inches(2.95), height=Inches(2.22))
    _add_planned_panel(slide, plans[PROMPT], _field_text(units, PROMPT), x=Inches(1.05), y=Inches(5.40), width=Inches(11.15), height=Inches(1.18))


def _render_controls(slide, units: tuple[DiscussionUnit, ...]):
    table = add_three_line_table(slide, ("变量", "固定/改变", "审查点"), (
        ("数据与 split", "固定", "相同受试者边界"), ("时序主干与输出", "固定", "PatchMixer / THO soft-z"),
        ("训练与评价", "固定", "相同损失与指标"), ("时频辅助表示", "唯一改变", "time-only / F0 / wide / bandenergy"),
    ), x=Inches(0.75), y=Inches(2.10), width=Inches(11.80), height=Inches(3.55), font_size=16)
    table._graphic_frame.name = "四对照变量表"
    add_evidence_boundary(slide, _field_text(units, LIMITS), x=Inches(1.05), y=Inches(5.88), width=Inches(11.15), height=Inches(0.72))


def _render_stft_table(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    table = add_three_line_table(slide, ("表示", "输入 shape", "保留信息", "tradeoff"), (
        ("F0 STFT", "B×239×37", "较细频率网格", "5 秒更新"),
        ("wide STFT", "B×160×73", "完整 0.05–8 Hz 谱图", "频率分辨力下降"),
        ("bandenergy", "B×5×73", "重叠频带强弱", "丢失带内峰位置"),
    ), x=Inches(0.70), y=Inches(2.05), width=Inches(11.90), height=Inches(3.25), font_size=16)
    table._graphic_frame.name = "STFT shape tradeoff 表"
    prompt = next(panel for panel in spec.panel_plan if panel.role == "discussion")
    _add_planned_panel(slide, prompt, _field_text(units, PROMPT), x=Inches(1.05), y=Inches(5.52), width=Inches(11.15), height=Inches(1.05))


def _canonical_rows(repo_root: Path):
    return load_canonical_overall_metrics(repo_root)


def _render_canonical(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...], repo_root: Path, *, secondary: bool):
    rows = _canonical_rows(repo_root)
    if secondary:
        headers = ("方案", "relative envelope corr", "local RR MAE bpm")
        values = tuple((MODEL_NAMES[row["label"]], f"{row['relative_envelope_corr']:.6f}", f"{row['local_rr_mae_bpm']:.6f}") for row in rows)
    else:
        headers = ("方案", "robust RR MAE bpm", "count bpm MAE", "lag4 corr（无量纲）")
        values = tuple((MODEL_NAMES[row["label"]], f"{row['robust_rr_mae_bpm']:.6f}", f"{row['count_bpm_mae']:.6f}", f"{row['lag4_corr']:.6f}") for row in rows)
    table = add_three_line_table(slide, headers, values, x=Inches(0.62), y=Inches(2.05), width=Inches(12.05), height=Inches(3.35), font_size=16)
    table._graphic_frame.name = "四方案核心指标（二）" if secondary else "四方案核心指标（一）"
    if secondary:
        add_discussion_box(slide, _field_text(units, PROMPT), x=Inches(1.05), y=Inches(5.72), width=Inches(11.15), height=Inches(1.05))
    else:
        add_evidence_boundary(slide, _field_text(units, EVIDENCE) + "\n" + _field_text(units, LIMITS), x=Inches(1.05), y=Inches(5.45), width=Inches(11.15), height=Inches(1.35))


CASE_LABELS = {"case_row_640": "四模型均纠正", "case_row_873": "四模型均失败", "case_row_1353": "模型间分歧", "case_row_3584": "阈值边界"}


def _case_panel(slide, *, title: str, body: str, x, y, width, height, accent, name: str):
    background = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    background.name = f"{name}背景"
    background.fill.solid(); background.fill.fore_color.rgb = WHITE
    background.line.color.rgb = accent
    add_text_box(slide, title, x=x + Inches(0.12), y=y + Inches(0.07), width=width - Inches(0.24), height=Inches(0.36), font_size=18, bold=True, name=f"{name}标题", align=PP_ALIGN.CENTER)
    add_text_box(slide, body, x=x + Inches(0.18), y=y + Inches(0.46), width=width - Inches(0.36), height=height - Inches(0.54), font_size=18, name=f"{name}正文", align=PP_ALIGN.LEFT)


def _render_case(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...], path: Path):
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.50))
    background.name = "案例页白色背景"
    background.fill.solid(); background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    tree = slide.shapes._spTree
    tree.remove(background._element)
    tree.insert(2, background._element)
    title = next(shape for shape in slide.shapes if shape.name == "页面标题")
    title.width = Inches(11.18)
    add_text_box(slide, f"类别｜{CASE_LABELS[spec.asset_key or '']}；本页按输入、目标、预测、频谱与指标完整回读。", x=Inches(0.80), y=Inches(1.78), width=Inches(11.15), height=Inches(0.42), font_size=18, name="副标题", align=PP_ALIGN.CENTER)
    _picture(slide, path, x=Inches(0.62), y=Inches(2.28), width=Inches(7.05), height=Inches(3.78))
    _case_panel(slide, title="证据边界", body="单窗案例不能说明总体效应；类别与选择规则须先冻结。", x=Inches(7.95), y=Inches(2.28), width=Inches(4.30), height=Inches(1.22), accent=SECONDARY_ORANGE, name="证据边界")
    _case_panel(slide, title="完整复核", body="• 固定类别与 row。\n• 同一 seed 四模型。\n• 回读 BCG/THO、频谱与指标。", x=Inches(7.95), y=Inches(3.68), width=Inches(4.30), height=Inches(1.82), accent=SCNU_BLUE, name="完整复核")
    _case_panel(slide, title="需要建议 / 待决策", body="还需哪些反例类别来暴露 gate 风险？", x=Inches(7.95), y=Inches(5.55), width=Inches(4.30), height=Inches(1.30), accent=SECONDARY_ORANGE, name="讨论框")


def _render_minimal_next(slide, spec: SlideSpec, units: tuple[DiscussionUnit, ...]):
    plans = {panel.field: panel for panel in spec.panel_plan}
    _add_planned_panel(slide, plans[METHOD], _field_text(units, METHOD), x=Inches(0.72), y=Inches(2.02), width=Inches(6.15), height=Inches(2.72))
    _add_planned_panel(slide, plans[RATIONALE], _field_text(units, RATIONALE), x=Inches(7.10), y=Inches(2.02), width=Inches(2.55), height=Inches(2.72))
    _add_planned_panel(slide, plans[LIMITS], _field_text(units, LIMITS), x=Inches(9.88), y=Inches(2.02), width=Inches(2.55), height=Inches(2.72))
    _add_planned_panel(slide, plans[PROMPT], _field_text(units, PROMPT), x=Inches(1.05), y=Inches(5.25), width=Inches(11.15), height=Inches(1.25))


def build_section_slide(prs: Presentation, section: str, number: int, page: int):
    slide = new_content_slide(prs, f"{number:02d}｜{section}", page_number=page)
    add_text_box(slide, "本章沿“问题 → 步骤 → 参数/shape → 证据 → 边界 → 决策”展开。", x=Inches(1.15), y=Inches(2.55), width=Inches(11.05), height=Inches(1.05), font_size=25, bold=True, name="章节导航", align=PP_ALIGN.CENTER)
    add_text_box(slide, "章节导航", x=Inches(4.65), y=Inches(4.35), width=Inches(4.05), height=Inches(0.70), font_size=20, color=SCNU_BLUE, bold=True, name="章节标签", align=PP_ALIGN.CENTER)
    return slide


def build_discussion_slide(prs: Presentation, spec: SlideSpec, page: int, *, repo_root: Path, assets: Mapping[str, Path]):
    units = _units(spec)
    slide = _base(prs, spec, page, units)
    path = assets.get(spec.asset_key or "")
    if spec.builder in {"asset", "softz", "asset_formula"}:
        if path is None:
            raise FileNotFoundError(f"缺少 SlideSpec {spec.key} 的真实资产")
        _render_asset(slide, spec, units, path)
        if spec.builder == "asset_formula":
            add_text_box(slide, FORMULAS[spec.key], x=Inches(0.84), y=Inches(1.78), width=Inches(11.40), height=Inches(0.30), font_size=18, bold=True, name="公式", align=PP_ALIGN.CENTER)
    elif spec.builder == "generic" or spec.builder == "decision_pair":
        _render_generic(slide, spec, units)
    elif spec.builder == "provenance":
        _render_provenance(slide, spec, units)
    elif spec.builder == "sample":
        _render_sample(slide, spec, units)
    elif spec.builder == "target_pipeline":
        _render_target_pipeline(slide, spec, units)
    elif spec.builder == "loss_goal":
        _render_loss_goal(slide, spec, units)
    elif spec.builder == "loss_weights":
        _render_loss_weights(slide, spec, units)
    elif spec.builder == "checkpoint_flow":
        _render_checkpoint_flow(slide, spec, units)
    elif spec.builder == "evidence_page":
        _render_evidence_page(slide, spec, units)
    elif spec.builder == "local_rr":
        if path is None:
            raise FileNotFoundError("缺少 local RR 真实资产")
        _render_local_rr(slide, spec, units, path)
    elif spec.builder == "formula":
        _render_formula(slide, spec, units)
    elif spec.builder == "controls":
        _render_controls(slide, units)
    elif spec.builder == "stft_table":
        _render_stft_table(slide, spec, units)
    elif spec.builder == "canonical_primary":
        _render_canonical(slide, spec, units, repo_root, secondary=False)
    elif spec.builder == "canonical_secondary":
        _render_canonical(slide, spec, units, repo_root, secondary=True)
    elif spec.builder == "case":
        if path is None:
            raise FileNotFoundError(f"缺少案例资产 {spec.asset_key}")
        _render_case(slide, spec, units, path)
    elif spec.builder == "minimal_next":
        _render_minimal_next(slide, spec, units)
    else:
        raise ValueError(f"未知 discussion builder: {spec.builder}")
    return [slide]
