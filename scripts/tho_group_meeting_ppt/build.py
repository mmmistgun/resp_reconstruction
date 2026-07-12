from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .charts import MODEL_NAMES, MODEL_ORDER, ChartData, build_all_charts, load_chart_data
from .content import BACKUP_SLIDES, MAIN_SLIDES, SlideSpec
from .theme import (
    BODY_BLACK,
    NOTE_GRAY,
    SCNU_BLUE,
    SECONDARY_ORANGE,
    WHITE,
    add_body_text,
    add_card,
    add_placeholder,
    add_source_note,
    add_takeaway,
    add_text_box,
    add_three_line_table,
    new_content_slide,
)


REPORT_DIR = Path("docs/stage_reports/20260708")
ASSET_DIR = REPORT_DIR / "generated_assets"
CASE_DIR = Path("runs/bcg_second_harmonic_20260710/figures")


def remove_all_sample_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def _keep_template_title_only(prs: Presentation) -> None:
    """保留模板第一张及其母版关系，只移除示例正文。"""
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids)[1:]:
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def _prepare_discussion_title(prs: Presentation) -> None:
    slide = prs.slides[0]
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    title = "THO research v2｜从整晚 BCG 到呼吸重建"
    subtitle = "全流程研究方法细节讨论｜证据、边界与下一轮决策"
    for shape, text, size in zip(text_shapes[:2], (title, subtitle), (30, 20), strict=False):
        shape.text_frame.clear()
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run(); run.text = text
        run.font.name = "Microsoft YaHei"; run.font.size = Pt(size)
        run.font.bold = size == 30; run.font.color.rgb = BODY_BLACK
    if len(text_shapes) >= 2:
        text_shapes[1].left = Inches(2.00)
        text_shapes[1].top = Inches(5.45)
        text_shapes[1].width = Inches(9.33)
        text_shapes[1].height = Inches(0.70)
    if len(text_shapes) < 2:
        add_text_box(slide, subtitle, x=Inches(2.0), y=Inches(4.8), width=Inches(9.3), height=Inches(0.8), font_size=20, name="标题页副标题", align=PP_ALIGN.CENTER)


def _add_contained_picture(slide, path: Path, *, x, y, width, height, name: str = "结果图"):
    with Image.open(path) as image:
        image_width, image_height = image.size
    scale = min(width / image_width, height / image_height)
    draw_width = int(image_width * scale)
    draw_height = int(image_height * scale)
    draw_x = int(x + (width - draw_width) / 2)
    draw_y = int(y + (height - draw_height) / 2)
    picture = slide.shapes.add_picture(str(path), draw_x, draw_y, draw_width, draw_height)
    picture.name = name
    return picture


def _add_flow_node(slide, text: str, *, x, y, width, height, accent=SCNU_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.name = f"流程节点：{text}"
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.4)
    add_text_box(
        slide,
        text,
        x=x + Pt(6),
        y=y + Pt(4),
        width=width - Pt(12),
        height=height - Pt(8),
        font_size=15,
        bold=True,
        name=f"流程节点文字：{text}",
        align=PP_ALIGN.CENTER,
    )
    return shape


def _add_arrow(slide, *, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.name = "流程连接线"
    line.line.color.rgb = NOTE_GRAY
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True
    return line


def _title_slide(prs: Presentation, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""
    blue = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), prs.slide_width, Inches(2.65))
    blue.name = "模板深蓝色主视觉"
    blue.fill.solid()
    blue.fill.fore_color.rgb = SCNU_BLUE
    blue.line.fill.background()
    add_text_box(
        slide,
        spec.title,
        x=Inches(1.25),
        y=Inches(2.55),
        width=Inches(10.85),
        height=Inches(0.70),
        font_size=34,
        color=WHITE,
        bold=True,
        name="标题页标题",
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        spec.takeaway,
        x=Inches(1.30),
        y=Inches(3.38),
        width=Inches(10.75),
        height=Inches(0.62),
        font_size=20,
        color=WHITE,
        name="标题页副标题",
        align=PP_ALIGN.CENTER,
    )
    add_placeholder(
        slide,
        spec.placeholder or "【待补：汇报人、汇报日期】",
        x=Inches(4.25),
        y=Inches(5.55),
        width=Inches(4.85),
        height=Inches(0.55),
    )
    return slide


def _base_slide(prs: Presentation, spec: SlideSpec, page_number: int):
    slide = new_content_slide(prs, spec.title, page_number=page_number)
    add_takeaway(slide, spec.takeaway)
    add_source_note(slide, spec.sources)
    return slide


def _add_bullet_layout(slide, spec: SlideSpec, *, x=0.92, y=2.05, width=11.4, height=4.55, font_size=20):
    if spec.bullets:
        add_body_text(
            slide,
            spec.bullets,
            x=Inches(x),
            y=Inches(y),
            width=Inches(width),
            height=Inches(height),
            font_size=font_size,
            bullet=True,
        )


def _build_three_questions(slide, spec: SlideSpec) -> None:
    for index, (title, body) in enumerate(
        zip(("01", "02", "03"), spec.bullets, strict=True)
    ):
        add_card(
            slide,
            title,
            body,
            x=Inches(0.88 + 4.05 * index),
            y=Inches(2.15),
            width=Inches(3.55),
            height=Inches(3.25),
            accent=(SCNU_BLUE if index < 2 else SECONDARY_ORANGE),
        )


def _build_task_difficulties(slide, spec: SlideSpec) -> None:
    add_body_text(slide, spec.bullets, x=Inches(0.86), y=Inches(2.08), width=Inches(5.2), height=Inches(3.95), font_size=18, bullet=True)
    add_placeholder(slide, spec.placeholder or "【待补：典型波形】", x=Inches(6.35), y=Inches(2.08), width=Inches(5.95), height=Inches(3.95))


def _build_dataset(slide) -> None:
    rows = [
        ("训练集", "32", "10,141"),
        ("验证集", "7", "2,675"),
        ("测试集", "8", "2,310"),
    ]
    add_three_line_table(slide, ("划分", "受试者数", "可用窗口数"), rows, x=Inches(0.92), y=Inches(2.15), width=Inches(6.25), height=Inches(2.65), font_size=18)
    add_card(slide, "独立性", "训练、验证、测试受试者编号不重叠。", x=Inches(7.55), y=Inches(2.10), width=Inches(4.60), height=Inches(1.20))
    add_card(slide, "统计边界", "窗口长度 180 秒；同一受试者内相邻窗口重叠，且窗口数不均衡。", x=Inches(7.55), y=Inches(3.55), width=Inches(4.60), height=Inches(1.65), accent=SECONDARY_ORANGE)


def _build_preprocessing(slide, spec: SlideSpec) -> None:
    nodes = ("宽频 BCG\n0.03–20 Hz", "状态对齐", "分段 robust-z", "soft-z 压缩", "模型输入")
    for index, text in enumerate(nodes):
        x = Inches(0.55 + index * 2.55)
        _add_flow_node(slide, text, x=x, y=Inches(2.45), width=Inches(2.05), height=Inches(1.15), accent=(SECONDARY_ORANGE if index == 3 else SCNU_BLUE))
        if index < len(nodes) - 1:
            _add_arrow(slide, x1=x + Inches(2.05), y1=Inches(3.02), x2=x + Inches(2.50), y2=Inches(3.02))
    add_body_text(slide, ("THO 目标同样采用分段 robust-z / soft-z。", "训练读取导出字段，不在窗口读取时重新滤波或归一化。"), x=Inches(1.10), y=Inches(4.20), width=Inches(11.0), height=Inches(1.45), font_size=18, bullet=True)


def _build_model_flow(slide) -> None:
    nodes = (
        ("180 秒 BCG soft-z", 0.45, 2.15, SCNU_BLUE),
        ("时序分块", 2.90, 2.15, SCNU_BLUE),
        ("patch_mixer1d", 5.35, 2.15, SCNU_BLUE),
        ("时序 token", 7.80, 2.15, SCNU_BLUE),
        ("同一 BCG", 0.45, 4.30, SECONDARY_ORANGE),
        ("STFT / 频带能量", 2.90, 4.30, SECONDARY_ORANGE),
        ("时频 token", 5.35, 4.30, SECONDARY_ORANGE),
        ("pre-mixer 融合", 8.15, 3.25, SCNU_BLUE),
        ("THO-like soft-z", 10.55, 3.25, SCNU_BLUE),
    )
    for text, x, y, accent in nodes:
        _add_flow_node(slide, text, x=Inches(x), y=Inches(y), width=Inches(1.95), height=Inches(0.90), accent=accent)
    for x in (2.40, 4.85, 7.30):
        _add_arrow(slide, x1=Inches(x), y1=Inches(2.60), x2=Inches(x + 0.45), y2=Inches(2.60))
    for x in (2.40, 4.85):
        _add_arrow(slide, x1=Inches(x), y1=Inches(4.75), x2=Inches(x + 0.45), y2=Inches(4.75))
    _add_arrow(slide, x1=Inches(7.30), y1=Inches(4.75), x2=Inches(8.10), y2=Inches(3.78))
    _add_arrow(slide, x1=Inches(9.75), y1=Inches(3.70), x2=Inches(10.50), y2=Inches(3.70))


def _build_training(slide, spec: SlideSpec) -> None:
    add_body_text(slide, spec.bullets, x=Inches(0.88), y=Inches(2.05), width=Inches(5.25), height=Inches(3.95), font_size=17, bullet=True)
    rows = [
        ("包络", "整体呼吸强弱"),
        ("频谱", "呼吸频带分布"),
        ("平滑 / 高频", "抑制尖锐抖动"),
        ("相对包络", "局部增强 / 减弱"),
        ("signed corr / cos", "低频形态与方向"),
    ]
    add_three_line_table(slide, ("监督分项", "约束性质"), rows, x=Inches(6.50), y=Inches(2.05), width=Inches(5.65), height=Inches(3.75), font_size=15)


def _build_four_way(slide) -> None:
    rows = [
        ("G0_time_only", "无", "—", "—", "纯时序参照"),
        ("G0_f0_native", "宽频 STFT", "30 s / 5 s", "pre-mixer", "上一版参照"),
        ("G3_C_wide_8p0", "宽频 STFT", "20 s / 2.5 s", "pre-mixer", "当前基准候选"),
        ("G3_C_bandenergy", "5 频带能量", "20 s / 2.5 s", "pre-mixer", "替代或条件修正"),
    ]
    add_three_line_table(slide, ("方案", "时频表示", "窗 / 步长", "融合位置", "回答的问题"), rows, x=Inches(0.55), y=Inches(2.12), width=Inches(12.15), height=Inches(3.75), font_size=13)


def _build_metrics(slide) -> None:
    rows = [
        ("稳健带通 RR 误差", "整窗主要呼吸节律", "越低越好"),
        ("周期计数 bpm 误差", "180 秒内周期数量", "越低越好"),
        ("4 秒时延校正相关", "允许时延后的低频形态", "越高越好"),
        ("相对包络相关", "局部呼吸增强 / 减弱", "越高越好"),
        ("局部 RR MAE / 相关", "滑窗节律变化", "低 / 高更好"),
    ]
    add_three_line_table(slide, ("指标", "回答的问题", "读表方向"), rows, x=Inches(0.75), y=Inches(2.08), width=Inches(11.85), height=Inches(3.85), font_size=15)


def _overall_rows(data: ChartData):
    return [
        (
            MODEL_NAMES[model],
            f"{data.overall.loc[model, 'robust_rr']:.3f}",
            f"{data.overall.loc[model, 'count_bpm']:.3f}",
            f"{data.overall.loc[model, 'lag4_corr']:.3f}",
            f"{data.overall.loc[model, 'local_rr_mae']:.3f}",
            f"{data.overall.loc[model, 'local_rr_corr']:.3f}",
        )
        for model in MODEL_ORDER
    ]


def _build_overall(slide, data: ChartData, assets: dict[str, Path]) -> None:
    del assets
    bold = {(3, 1), (4, 2), (3, 3), (4, 4), (4, 5)}
    add_three_line_table(
        slide,
        ("方案", "稳健 RR", "计数 bpm", "lag4 corr", "local RR MAE", "local RR corr"),
        _overall_rows(data),
        x=Inches(0.50),
        y=Inches(2.18),
        width=Inches(12.30),
        height=Inches(3.55),
        bold_cells=bold,
        font_size=16,
    )
    add_text_box(
        slide,
        "wide：稳健 RR 与时延校正形态更优；bandenergy：周期计数与局部 RR 略优。",
        x=Inches(1.10),
        y=Inches(5.95),
        width=Inches(11.05),
        height=Inches(0.48),
        font_size=17,
        bold=True,
        name="正文结果解读",
        align=PP_ALIGN.CENTER,
    )


def _build_evidence_slide(slide, spec: SlideSpec, assets: dict[str, Path], kind: str) -> None:
    if kind == "stft":
        _add_contained_picture(slide, assets["overall_metrics"], x=Inches(0.72), y=Inches(2.05), width=Inches(7.1), height=Inches(4.35))
        add_body_text(slide, spec.bullets, x=Inches(8.10), y=Inches(2.25), width=Inches(4.25), height=Inches(3.75), font_size=17, bullet=True)
    elif kind == "tail":
        _add_contained_picture(slide, assets["tail_metrics"], x=Inches(0.75), y=Inches(2.00), width=Inches(11.8), height=Inches(4.75))
    else:
        _add_bullet_layout(slide, spec, y=2.15, height=3.8, font_size=18)


def _build_scope(slide, spec: SlideSpec) -> None:
    for index, (title, body) in enumerate((
        ("能回答", "收益主要出现在哪些已知窗口？"),
        ("不能回答", "推理时应对哪些窗口开启修正？"),
        ("下一步", "构造完全由 BCG 输入计算的无泄漏判据。"),
    )):
        add_card(slide, title, body, x=Inches(0.75 + 4.13 * index), y=Inches(2.25), width=Inches(3.72), height=Inches(2.85), accent=(SCNU_BLUE if index != 1 else SECONDARY_ORANGE))


def _build_strata(slide, spec: SlideSpec, assets: dict[str, Path]) -> None:
    _add_contained_picture(slide, assets["stratified_roles"], x=Inches(0.70), y=Inches(1.95), width=Inches(8.45), height=Inches(4.90))
    add_card(slide, "解读", spec.takeaway, x=Inches(9.30), y=Inches(2.45), width=Inches(3.10), height=Inches(2.25), accent=SECONDARY_ORANGE)


def _build_count_error(slide) -> None:
    rows = [
        ("count-error", "纯时序", "1.239", "1.428", "1.185"),
        ("count-error", "宽频 STFT", "1.145", "1.260", "1.171"),
        ("count-error", "频带能量", "1.171", "1.201", "1.157"),
        ("clean-count", "纯时序", "0.216", "0.000", "0.310"),
        ("clean-count", "宽频 STFT", "0.193", "0.066", "0.312"),
        ("clean-count", "频带能量", "0.202", "0.068", "0.314"),
    ]
    add_three_line_table(slide, ("分层", "方案", "稳健 RR", "计数 bpm", "local RR MAE"), rows, x=Inches(0.70), y=Inches(2.05), width=Inches(11.90), height=Inches(3.60), bold_cells={(2, 2), (3, 3), (3, 4), (5, 2), (4, 3), (4, 4)}, font_size=14)
    add_text_box(slide, "clean-count 中，纯时序计数误差本来为 0；时频方案会引入少量非零误差。", x=Inches(1.15), y=Inches(5.85), width=Inches(11.0), height=Inches(0.55), font_size=16, bold=True, name="正文提示", align=PP_ALIGN.CENTER)


def _build_harmonic_mechanism(slide, spec: SlideSpec) -> None:
    nodes = ("THO 基频 f", "BCG 基频附近能量", "BCG 主峰接近 2f", "网络输出回到 f")
    accents = (SCNU_BLUE, SCNU_BLUE, SECONDARY_ORANGE, SCNU_BLUE)
    for index, (text, accent) in enumerate(zip(nodes, accents, strict=True)):
        x = Inches(0.65 + index * 3.10)
        _add_flow_node(slide, text, x=x, y=Inches(2.30), width=Inches(2.45), height=Inches(1.15), accent=accent)
        if index < 3:
            _add_arrow(slide, x1=x + Inches(2.45), y1=Inches(2.88), x2=x + Inches(3.05), y2=Inches(2.88))
    add_body_text(slide, spec.bullets, x=Inches(1.00), y=Inches(4.15), width=Inches(11.25), height=Inches(1.60), font_size=17, bullet=True)


def _build_harmonic_coverage(slide, assets: dict[str, Path]) -> None:
    _add_contained_picture(slide, assets["harmonic_coverage"], x=Inches(0.55), y=Inches(2.00), width=Inches(7.15), height=Inches(4.40))
    add_card(slide, "可判定窗口", "1,924 / 2,310\n占全部测试窗口 83.29%", x=Inches(7.90), y=Inches(2.10), width=Inches(4.15), height=Inches(1.15))
    add_card(slide, "阳性并集", "452 个窗口｜19.57%\n覆盖 6 名受试者", x=Inches(7.90), y=Inches(3.55), width=Inches(4.15), height=Inches(1.15), accent=SECONDARY_ORANGE)
    add_card(slide, "证据集中性", "strong 214 个窗口；阳性与 strong 均集中在少数受试者。", x=Inches(7.90), y=Inches(5.00), width=Inches(4.15), height=Inches(1.25), accent=SECONDARY_ORANGE)


def _build_harmonic_results(slide, data: ChartData, assets: dict[str, Path]) -> None:
    del data
    _add_contained_picture(slide, assets["harmonic_model_results"], x=Inches(0.60), y=Inches(2.00), width=Inches(12.10), height=Inches(4.35))


def _build_balanced_cases(slide, repo_root: Path) -> None:
    cases = (
        ("四模型均纠正", repo_root / CASE_DIR / "all_corrected_seed_20260837_row_640.png", "row 640"),
        ("四模型均失败", repo_root / CASE_DIR / "all_not_corrected_seed_20260837_row_873.png", "row 873"),
        ("模型间分歧", repo_root / CASE_DIR / "model_disagreement_seed_20260837_row_1353.png", "row 1353"),
        ("阈值附近", repo_root / CASE_DIR / "threshold_boundary_seed_20260837_row_3584.png", "row 3584"),
    )
    for index, (title, path, row) in enumerate(cases):
        column = index % 2
        line = index // 2
        x = Inches(0.60 + 6.20 * column)
        y = Inches(2.00 + 2.35 * line)
        add_text_box(slide, f"{title}｜{row}", x=x, y=y, width=Inches(5.80), height=Inches(0.30), font_size=13, bold=True, name="案例标题", align=PP_ALIGN.CENTER)
        if path.exists():
            _add_contained_picture(slide, path, x=x, y=y + Inches(0.35), width=Inches(5.80), height=Inches(1.82), name=f"案例图：{row}")
        else:
            add_placeholder(slide, f"【待补：{title}；建议来源：{path}】", x=x, y=y + Inches(0.35), width=Inches(5.80), height=Inches(1.82))


def _build_decision(slide) -> None:
    cards = (
        ("当前基准", "G3_C_wide_8p0\n宽频、高时间分辨率 STFT", SCNU_BLUE),
        ("条件修正候选", "G3_C_bandenergy\n只在难恢复窗口探索", SECONDARY_ORANGE),
        ("停止动作", "不补评非重点分支\n不继续无差别扩大输入", NOTE_GRAY),
    )
    for index, (title, body, accent) in enumerate(cards):
        add_card(slide, title, body, x=Inches(0.80 + 4.15 * index), y=Inches(2.20), width=Inches(3.72), height=Inches(2.90), accent=accent)
    add_text_box(slide, "下一阶段研究问题：什么时候使用辅助时频信息？", x=Inches(2.10), y=Inches(5.50), width=Inches(9.15), height=Inches(0.60), font_size=22, color=SCNU_BLUE, bold=True, name="强调问题", align=PP_ALIGN.CENTER)


def _build_takeaways(slide, spec: SlideSpec) -> None:
    for index, text in enumerate(("STFT 值得保留", "wide 是当前基准", "bandenergy 用于条件修正探索")):
        add_card(slide, f"0{index + 1}", text, x=Inches(0.85 + 4.10 * index), y=Inches(2.00), width=Inches(3.60), height=Inches(1.55), accent=(SECONDARY_ORANGE if index == 2 else SCNU_BLUE))
    add_body_text(slide, spec.bullets, x=Inches(1.05), y=Inches(4.05), width=Inches(11.10), height=Inches(2.15), font_size=17, bullet=True)


def build_main_slide(prs: Presentation, spec: SlideSpec, page_number: int, *, repo_root: Path, data: ChartData, assets: dict[str, Path]):
    if spec.key == "title":
        return _title_slide(prs, spec)
    slide = _base_slide(prs, spec, page_number)
    if spec.key == "three_questions":
        _build_three_questions(slide, spec)
    elif spec.key == "task_and_difficulties":
        _build_task_difficulties(slide, spec)
    elif spec.key == "dataset_split":
        _build_dataset(slide)
    elif spec.key == "preprocessing":
        _build_preprocessing(slide, spec)
    elif spec.key == "model_flow":
        _build_model_flow(slide)
    elif spec.key == "training_supervision":
        _build_training(slide, spec)
    elif spec.key == "four_way_control":
        _build_four_way(slide)
    elif spec.key == "metric_framework":
        _build_metrics(slide)
    elif spec.key == "overall_test_results":
        _build_overall(slide, data, assets)
    elif spec.key == "stft_adds_information":
        _build_evidence_slide(slide, spec, assets, "stft")
    elif spec.key == "tail_errors":
        _build_evidence_slide(slide, spec, assets, "tail")
    elif spec.key == "stratification_scope":
        _build_scope(slide, spec)
    elif spec.key in {"hard_easy_low_spectrum", "rr_roles"}:
        _build_strata(slide, spec, assets)
    elif spec.key == "count_error":
        _build_count_error(slide)
    elif spec.key == "harmonic_motivation":
        _build_harmonic_mechanism(slide, spec)
    elif spec.key == "harmonic_coverage":
        _build_harmonic_coverage(slide, assets)
    elif spec.key == "harmonic_model_results":
        _build_harmonic_results(slide, data, assets)
    elif spec.key == "balanced_cases":
        _build_balanced_cases(slide, repo_root)
    elif spec.key == "stage_decision":
        _build_decision(slide)
    elif spec.key == "next_stage_takeaways":
        _build_takeaways(slide, spec)
    else:
        _add_bullet_layout(slide, spec)
        if spec.placeholder:
            add_placeholder(slide, spec.placeholder, x=Inches(6.40), y=Inches(3.75), width=Inches(5.65), height=Inches(1.60))
    return slide


def build_presentation(
    *,
    template: Path,
    output: Path,
    include_backup: bool = True,
    repo_root: Path | None = None,
) -> Path:
    repo_root = repo_root or Path(__file__).resolve().parents[2]
    assets = build_all_charts(repo_root, repo_root / ASSET_DIR)
    data = load_chart_data(repo_root)
    prs = Presentation(template)
    remove_all_sample_slides(prs)
    for page_number, spec in enumerate(MAIN_SLIDES, start=1):
        build_main_slide(prs, spec, page_number, repo_root=repo_root, data=data, assets=assets)
    if include_backup:
        from .backup import build_backup_slides

        build_backup_slides(prs, repo_root=repo_root, data=data, assets=assets, start_page=len(MAIN_SLIDES) + 1)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


def build_discussion_assets(repo_root: Path, output_dir: Path | None = None) -> dict[str, Path]:
    """按任务 3–6 顺序刷新真实位图资产，并合并返回路径。"""
    from .case_figures import build_case_assets
    from .metric_figures import build_metric_assets
    from .model_figures import build_model_assets
    from .signal_figures import build_signal_assets

    output_dir = output_dir or repo_root / REPORT_DIR / "generated_assets/discussion"
    combined: dict[str, Path] = {}
    for builder in (build_signal_assets, build_model_assets, build_metric_assets, build_case_assets):
        assets, _ = builder(repo_root, output_dir)
        combined.update(assets)
    return combined


def build_discussion_presentation(
    *,
    template: Path,
    output: Path,
    repo_root: Path | None = None,
) -> Path:
    """组装首次接触项目也可完整阅读的研究讨论型正文。"""
    from .detail_slides import SLIDE_GROUPS, build_discussion_slide, build_section_slide, resolve_discussion_assets

    root = (repo_root or Path(__file__).resolve().parents[2]).resolve()
    prs = Presentation(template)
    _keep_template_title_only(prs)
    _prepare_discussion_title(prs)
    assets = resolve_discussion_assets(root)
    page = 2
    current_section = None
    section_number = 0
    for spec in SLIDE_GROUPS:
        if spec.section != current_section:
            section_number += 1
            build_section_slide(prs, spec.section, section_number, page)
            page += 1
            current_section = spec.section
        slides = build_discussion_slide(prs, spec, page, repo_root=root, assets=assets)
        page += len(slides)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    _sanitize_discussion_package(output)
    return output


def _sanitize_discussion_package(path: Path) -> None:
    """移除模板残留的隐藏示例文案，避免其留在交付包属性中。"""
    replacements = {
        "论文分享": "研究讨论",
        "汇报人：xxx": "THO research v2",
        "2021/12/21": "",
    }
    temporary = path.with_suffix(".sanitized.pptx")
    with ZipFile(path, "r") as source, ZipFile(temporary, "w") as target:
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename.endswith((".xml", ".rels")):
                text = data.decode("utf-8")
                for old, new in replacements.items():
                    text = text.replace(old, new)
                data = text.encode("utf-8")
            target.writestr(info, data)
    temporary.replace(path)
